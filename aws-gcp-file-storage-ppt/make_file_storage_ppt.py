#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""AWS vs GCP 文件存储对比 PPT (3组产品) - NetApp风格"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# 配色 (TOOLS.md NetApp风)
NETAPP_BLUE=RGBColor(0x00,0x67,0xC5)
ORANGE=RGBColor(0xF5,0x82,0x20)
GREEN=RGBColor(0x2E,0x9E,0x5B)
DARK=RGBColor(0x1A,0x2B,0x4A)
LIGHT_BG=RGBColor(0xF2,0xF6,0xFB)
GRAY=RGBColor(0x5A,0x6B,0x82)
AWS_ORANGE=RGBColor(0xFF,0x99,0x00)
GCP_BLUE=RGBColor(0x42,0x85,0xF4)
WHITE=RGBColor(0xFF,0xFF,0xFF)
FONT="Microsoft YaHei"

prs=Presentation()
prs.slide_width=Inches(13.333); prs.slide_height=Inches(7.5)
SW,SH=prs.slide_width,prs.slide_height
blank=prs.slide_layouts[6]

def solid(shape,color):
    shape.fill.solid(); shape.fill.fore_color.rgb=color; shape.line.fill.background()
def box(sl,l,t,w,h,color=None):
    from pptx.enum.shapes import MSO_SHAPE
    s=sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,l,t,w,h)
    if color: solid(s,color)
    else: s.fill.background(); s.line.fill.background()
    return s
def rect(sl,l,t,w,h,color):
    from pptx.enum.shapes import MSO_SHAPE
    s=sl.shapes.add_shape(MSO_SHAPE.RECTANGLE,l,t,w,h); solid(s,color); return s
def txt(sl,l,t,w,h,text,size=14,color=DARK,bold=False,align=PP_ALIGN.LEFT,anchor=MSO_ANCHOR.TOP,font=FONT):
    tb=sl.shapes.add_textbox(l,t,w,h); tf=tb.text_frame; tf.word_wrap=True
    tf.vertical_anchor=anchor
    lines=text.split("\n")
    for i,ln in enumerate(lines):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph()
        p.alignment=align
        r=p.add_run(); r.text=ln
        r.font.size=Pt(size); r.font.bold=bold; r.font.color.rgb=color; r.font.name=font
    return tb

def title_bar(sl,title,sub=""):
    rect(sl,0,0,SW,Inches(1.05),NETAPP_BLUE)
    rect(sl,0,Inches(1.05),SW,Pt(4),ORANGE)
    txt(sl,Inches(0.5),Inches(0.12),Inches(12),Inches(0.6),title,26,WHITE,True,anchor=MSO_ANCHOR.MIDDLE)
    if sub: txt(sl,Inches(0.5),Inches(0.62),Inches(12),Inches(0.35),sub,13,RGBColor(0xD5,0xE3,0xF2))

# ============ 封面 ============
s=prs.slides.add_slide(blank)
rect(s,0,0,SW,SH,DARK)
rect(s,0,0,Inches(0.35),SH,ORANGE)
txt(s,Inches(0.9),Inches(2.3),Inches(11.5),Inches(1.2),"AWS vs GCP 文件存储服务对比",40,WHITE,True)
txt(s,Inches(0.9),Inches(3.5),Inches(11.5),Inches(0.6),"EFS·Filestore | FSx Lustre·Managed Lustre | Parallelstore(DAOS) | FSx NetApp ONTAP·NetApp Volumes",16,RGBColor(0x9F,0xB8,0xD8))
txt(s,Inches(0.9),Inches(4.3),Inches(11),Inches(0.5),"基于两家官方文档整理 · 2026-08",13,GRAY)

# ============ 目录/概览 ============
s=prs.slides.add_slide(blank)
title_bar(s,"三组对比产品概览","按场景配对：基础NFS / 并行HPC / 企业级ONTAP")
rows=[
 ("场景","AWS","GCP"),
 ("基础托管 NFS","Amazon EFS","Filestore"),
 ("并行FS-Lustre系","FSx for Lustre","Managed Lustre (DDN)"),
 ("并行FS-对象系","(无直接对标)","Parallelstore (DAOS)"),
 ("企业级 ONTAP","FSx for NetApp ONTAP","NetApp Volumes"),
]
tw=Inches(12); th=Inches(3.2); l=Inches(0.66); t=Inches(1.7)
rh=th/len(rows); cw=tw/3
for ri,row in enumerate(rows):
    for ci,cell in enumerate(row):
        cellbg = NETAPP_BLUE if ri==0 else (LIGHT_BG if ri%2 else WHITE)
        r=rect(s,l+cw*ci,t+rh*ri,cw,rh,cellbg)
        col = WHITE if ri==0 else (AWS_ORANGE if ci==1 else (GCP_BLUE if ci==2 else DARK))
        bold = ri==0 or ci==0
        txt(s,l+cw*ci,t+rh*ri+Inches(0.12),cw,rh,cell,15 if ri==0 else 14,col,bold,PP_ALIGN.CENTER,MSO_ANCHOR.MIDDLE)
txt(s,Inches(0.66),Inches(5.3),Inches(12),Inches(1.2),
    "• 四组配对：EFS↔Filestore(NFS)；FSx Lustre↔Managed Lustre(都是Lustre)；Parallelstore(DAOS)是GCP独有的对象式并行FS；两个ONTAP同源NetApp\n• ⚠️ GCP有两个并行文件系统：Managed Lustre(Lustre内核) 和 Parallelstore(Intel DAOS内核)，是两个不同产品\n• 数据来源：AWS/GCP 官方文档 + GCP官方博客",
    12,GRAY)

# ============ 第1组: EFS vs Filestore ============
s=prs.slides.add_slide(blank)
title_bar(s,"① Amazon EFS  vs  GCP Filestore","基础托管 NFS 文件存储")
data=[
 ("维度","Amazon EFS","GCP Filestore"),
 ("协议","NFS v4.0/4.1","NFS v3"),
 ("层级/模式","Regional / One Zone;\n吞吐:Elastic/Provisioned/Bursting","Zonal / Regional / Enterprise(GKE)\n/ Basic HDD·SSD"),
 ("容量","无需预置,自动弹性伸缩(PB级)","Zonal 1-100TiB; Regional 100GiB-100TiB;\nBasic 最高63.9TiB"),
 ("性能(读吞吐)","Elastic: 20-60 GiBps/FS;\n单客户端 1500 MiBps","可配置(custom performance);\n随容量/层级"),
 ("最大IOPS","Elastic 读90万-250万/写5万\n(可申请提升10x)","随层级/容量"),
 ("可用性","Regional=多AZ; One Zone=单AZ","Zonal=单区; Regional=跨区抗区故障"),
 ("典型场景","无预置弹性共享/容器/通用NFS","HPC(Zonal)/区域高可用/GKE多share"),
]
tl=Inches(0.5); tt=Inches(1.55); tw=Inches(12.33); th=Inches(5.5)
rh=th/len(data); cw=[Inches(2.5),Inches(4.9),Inches(4.93)]
def draw_table(s,data,tl,tt,rh,cw):
    x=tl
    for ci in range(3):
        y=tt
        for ri,row in enumerate(data):
            bg=NETAPP_BLUE if ri==0 else (LIGHT_BG if ri%2 else WHITE)
            rect(s,x,y,cw[ci],rh,bg)
            if ri==0: col=WHITE
            elif ci==0: col=DARK
            elif ci==1: col=AWS_ORANGE
            else: col=GCP_BLUE
            txt(s,x+Inches(0.1),y,cw[ci]-Inches(0.15),rh,row[ci],
                12 if ri==0 else 10.5, col, ri==0 or ci==0, PP_ALIGN.LEFT if ci else PP_ALIGN.LEFT, MSO_ANCHOR.MIDDLE)
            y+=rh
        x+=cw[ci]
draw_table(s,data,tl,tt,rh,cw)

# ============ 第2组: FSx Lustre vs Managed Lustre ============
s=prs.slides.add_slide(blank)
title_bar(s,"② FSx for Lustre  vs  GCP Managed Lustre","并行文件系统 (HPC / AI 训练)")
data=[
 ("维度","FSx for Lustre","Managed Lustre (DDN)"),
 ("底层","AWS自研+开源Lustre","DDN EXAScaler 商业版Lustre"),
 ("吞吐档位","125/250/500/1000 MBps/TiB","125/250/500/1000 + Dynamic(25)"),
 ("吞吐性质","burst/baseline双档(会衰减)\n实测:125档@4.8TiB burst~1600→基线600","官方称 sustained 恒定吞吐"),
 ("对象存储集成","DRA + HSM 惰性加载\n(released/restore, 可释放空间)","GCS 批量全量传输(import/export)\n无HSM惰性加载"),
 ("元数据IOPS","可独立预置(1500-192000)","随容量+吞吐自动伸缩"),
 ("上限","多TBps","10 TBps / 80 PiB"),
 ("在线变配","支持,但OST断连~2min+6h间隔(实测)","增容量即提升"),
]
draw_table(s,data,tl,tt,rh,cw)

# ============ 第2组b: Parallelstore vs FSx Lustre ============
s=prs.slides.add_slide(blank)
title_bar(s,"②b GCP Parallelstore  vs  AWS FSx Lustre","并行文件系统 - 对象/DAOS 架构 (HPC/AI)")
data=[
 ("维度","AWS FSx for Lustre","GCP Parallelstore"),
 ("底层架构","Lustre (POSIX并行FS)","Intel DAOS (对象+KV, 用户态)"),
 ("元数据","集中式MDS/MDT(可多MDT)","全分布式元数据(无集中MDS)"),
 ("介质/IO路径","块设备+内核态LNet","持久内存(SCM)+NVMe, 用户态+RDMA"),
 ("最大规模","数百TB~PB / 多TBps","100 TiB / ~115 GiB/s (当前最大部署)"),
 ("IOPS/延迟","network IOPS/TiB; 亚毫秒","读~300万/写~100万; ~0.3ms (@100TiB)"),
 ("对象存储集成","DRA+HSM惰性加载","import/export API(批量~20GB/s)"),
 ("擅长场景","超大顺序吞吐(大文件)","小文件/元数据密集/AI训练(自测快3.9x)"),
]
draw_table(s,data,tl,tt,rh,cw)

# ============ 第3组: FSx ONTAP vs NetApp Volumes ============
s=prs.slides.add_slide(blank)
title_bar(s,"③ FSx for NetApp ONTAP  vs  GCP NetApp Volumes","企业级 ONTAP (同源 NetApp 技术)")
data=[
 ("维度","FSx for NetApp ONTAP","GCP NetApp Volumes"),
 ("底层","ONTAP on AWS","NetApp+Google 托管ONTAP"),
 ("协议","NFS/SMB/iSCSI/NVMe","NFS(v3/4.1/4.2)/SMB/iSCSI/NVMe\n+S3(ONTAP-mode)"),
 ("产品形态","单服务+部署类型(Single/Multi-AZ)","多service level:Flex/Standard\n/Premium/Extreme"),
 ("SnapMirror","支持(on-prem→云 迁移/DR)","支持(hybrid replication, 同能力)"),
 ("FlexGroup/大卷","支持FlexGroup","Large Capacity Volume(=FlexGroup)\n1PiB/12.5GiB·s/6端点"),
 ("多HA pair","Single-AZ 最多12 HA pair(用户可配)\n→72GB/s","HA不对用户暴露(托管黑盒)"),
 ("HA/冗余","Multi-AZ HA pair 自动failover","服务级HA + 可选多区冗余"),
]
draw_table(s,data,tl,tt,rh,cw)

# ============ 关键差异总结 ============
s=prs.slides.add_slide(blank)
title_bar(s,"关键差异 & 选型建议","核心分野")
pts=[
 ("EFS vs Filestore",GREEN,"EFS 免预置全弹性、NFSv4、单FS吞吐上限更高(60GiBps); Filestore 分层清晰(Zonal/Regional/Enterprise)。基础NFS两家都成熟。"),
 ("Lustre 并行FS",ORANGE,"对象存储集成是关键:AWS HSM惰性加载(数据>>容量省钱) vs GCP批量传输(全量装进,训练零抖动); AWS burst衰减/GCP恒定; AWS元数据可独立调。"),
 ("Parallelstore(DAOS)",RGBColor(0x8E,0x44,0xAD),"GCP独有的对象式并行FS(Intel DAOS,非Lustre)。全分布式元数据+用户态RDMA,主打小文件/元数据/AI训练(0.3ms/300万读IOPS)。AWS无直接对标(FSx Lustre最接近但架构不同)。"),
 ("NetApp ONTAP",NETAPP_BLUE,"同源→功能高度重合(快照/克隆/SnapMirror/FlexGroup); 最大区别:AWS开放多HA pair scale-out(最多12,用户可配) vs GCP HA托管黑盒。"),
]
y=Inches(1.55)
for name,color,desc in pts:
    rect(s,Inches(0.5),y,Inches(0.14),Inches(1.15),color)
    txt(s,Inches(0.72),y+Inches(0.03),Inches(3.0),Inches(1.1),name,14,color,True,anchor=MSO_ANCHOR.MIDDLE)
    box(s,Inches(3.85),y,Inches(9.0),Inches(1.2),LIGHT_BG)
    txt(s,Inches(4.05),y+Inches(0.08),Inches(8.7),Inches(1.05),desc,11.5,DARK,anchor=MSO_ANCHOR.MIDDLE)
    y+=Inches(1.35)
txt(s,Inches(0.5),Inches(7.05),Inches(12),Inches(0.35),"数据来源:AWS/GCP官方文档 2026-08 | 部分性能为实测,以实际region/配置为准",11,GRAY)

out="/home/ubuntu/.openclaw/workspace/aws_gcp_file_storage_compare.pptx"
prs.save(out)
print("saved",out,"pages:",len(prs.slides._sldIdLst))
