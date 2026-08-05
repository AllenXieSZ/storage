#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
AWS vs GCP 网络产品深度对比 PPT
数据来源：AWS/GCP 官方文档核实（每页脚注标来源），2026-08-05。
风格：伟伟认可的 NetApp 风模板（16:9），AWS 橙 / GCP 蓝双列对照。
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

AWS_ORANGE = RGBColor(0xFF, 0x99, 0x00)
GCP_BLUE = RGBColor(0x42, 0x85, 0xF4)
MAIN_BLUE = RGBColor(0x00, 0x67, 0xC5)
DARK = RGBColor(0x1A, 0x2B, 0x4A)
GRAY = RGBColor(0x5A, 0x6B, 0x82)
LIGHT_BG = RGBColor(0xF2, 0xF6, 0xFB)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
ORANGE = RGBColor(0xF5, 0x82, 0x20)
GREEN = RGBColor(0x2E, 0x9E, 0x5B)
RED = RGBColor(0xD0, 0x3A, 0x2E)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
blank = prs.slide_layouts[6]


def add_slide():
    return prs.slides.add_slide(blank)


def bg(slide, color=WHITE):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    s.fill.solid(); s.fill.fore_color.rgb = color
    s.line.fill.background(); s.shadow.inherit = False
    slide.shapes._spTree.remove(s._element)
    slide.shapes._spTree.insert(2, s._element)
    return s


def box(slide, l, t, w, h, fill=None, line=None, line_w=1.0, shape=MSO_SHAPE.RECTANGLE):
    s = slide.shapes.add_shape(shape, l, t, w, h)
    if fill is None:
        s.fill.background()
    else:
        s.fill.solid(); s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line; s.line.width = Pt(line_w)
    s.shadow.inherit = False
    return s


def text(slide, l, t, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, wrap=True):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = Pt(2); tf.margin_right = Pt(2)
    tf.margin_top = Pt(1); tf.margin_bottom = Pt(1)
    if isinstance(runs, str):
        runs = [(runs, 18, DARK, False)]
    first = True
    for item in runs:
        txt, sz, col, bold = item[0], item[1], item[2], item[3]
        sp = item[4] if len(item) > 4 else None
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = align
        if sp is not None:
            p.space_before = Pt(sp)
        r = p.add_run(); r.text = txt
        r.font.size = Pt(sz); r.font.bold = bold
        r.font.color.rgb = col
        r.font.name = "Microsoft YaHei"
    return tb


def title_bar(slide, title, sub=None):
    box(slide, 0, 0, SW, Inches(1.05), fill=MAIN_BLUE)
    box(slide, 0, Inches(1.05), SW, Pt(4), fill=ORANGE)
    text(slide, Inches(0.55), Inches(0.12), Inches(12), Inches(0.8),
         [(title, 25, WHITE, True)], anchor=MSO_ANCHOR.MIDDLE)
    if sub:
        text(slide, Inches(0.6), Inches(1.18), Inches(12.2), Inches(0.4),
             [(sub, 13, GRAY, False)])


def two_col_table(slide, rows, top=Inches(1.75),
                  col_labels=("AWS", "GCP"),
                  x0=Inches(0.5), label_w=Inches(2.6), col_w=Inches(5.05),
                  row_h=Inches(0.62), fs=11.5):
    x_aws = x0 + label_w
    x_gcp = x_aws + col_w
    box(slide, x0, top, label_w, Inches(0.5), fill=DARK)
    box(slide, x_aws, top, col_w, Inches(0.5), fill=AWS_ORANGE)
    box(slide, x_gcp, top, col_w, Inches(0.5), fill=GCP_BLUE)
    text(slide, x0, top, label_w, Inches(0.5), [("对比维度", 13, WHITE, True)],
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    text(slide, x_aws, top, col_w, Inches(0.5), [(col_labels[0], 13, WHITE, True)],
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    text(slide, x_gcp, top, col_w, Inches(0.5), [(col_labels[1], 13, WHITE, True)],
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    y = top + Inches(0.5)
    for i, (lab, a, g) in enumerate(rows):
        rb = LIGHT_BG if i % 2 == 0 else WHITE
        box(slide, x0, y, label_w, row_h, fill=RGBColor(0xE8, 0xEE, 0xF6), line=WHITE, line_w=0.5)
        box(slide, x_aws, y, col_w, row_h, fill=rb, line=WHITE, line_w=0.5)
        box(slide, x_gcp, y, col_w, row_h, fill=rb, line=WHITE, line_w=0.5)
        text(slide, x0 + Inches(0.08), y, label_w - Inches(0.12), row_h,
             [(lab, 11, MAIN_BLUE, True)], anchor=MSO_ANCHOR.MIDDLE)
        text(slide, x_aws + Inches(0.1), y, col_w - Inches(0.18), row_h,
             [(a, fs, DARK, False)], anchor=MSO_ANCHOR.MIDDLE)
        text(slide, x_gcp + Inches(0.1), y, col_w - Inches(0.18), row_h,
             [(g, fs, DARK, False)], anchor=MSO_ANCHOR.MIDDLE)
        y += row_h


def source_note(slide, txt):
    text(slide, Inches(0.55), Inches(7.05), Inches(12.4), Inches(0.35),
         [(txt, 9, GRAY, False)])


def sep_slide(title, subtitle=None):
    s = add_slide(); bg(s, DARK)
    box(s, 0, Inches(3.0), Inches(0.28), Inches(1.4), fill=ORANGE)
    text(s, Inches(0.7), Inches(3.05), Inches(12), Inches(1.0),
         [(title, 34, WHITE, True)], anchor=MSO_ANCHOR.MIDDLE)
    if subtitle:
        text(s, Inches(0.72), Inches(4.2), Inches(12), Inches(0.6),
             [(subtitle, 15, RGBColor(0x9F, 0xB3, 0xD1), False)])
    return s


# ============ Slide 1: 封面 ============
s = add_slide(); bg(s, DARK)
box(s, 0, 0, Inches(0.35), SH, fill=ORANGE)
box(s, Inches(0.9), Inches(2.5), Inches(6.0), Pt(3), fill=ORANGE)
text(s, Inches(0.9), Inches(1.7), Inches(11.5), Inches(1.4),
     [("AWS vs GCP", 46, WHITE, True), ("网络产品深度对比", 34, RGBColor(0xCF,0xDD,0xF0), True)])
text(s, Inches(0.9), Inches(3.9), Inches(11.5), Inches(1.2),
     [("外网 · 内网 · 加速 · CDN · DNS · 安全 · DDoS · 稳定性", 16, AWS_ORANGE, False),
      ("11 个维度 · 每条论点均基于官方文档核实", 13, GRAY, False, 8)])
text(s, Inches(0.9), Inches(6.7), Inches(11.5), Inches(0.5),
     [("数据来源：AWS/GCP 官方文档 + SLA/定价页  |  2026-08", 12, GRAY, False)])

# ============ Slide 2: 对标关系总览 ============
s = add_slide(); bg(s, WHITE)
title_bar(s, "网络产品对标总览", "8 大维度 · AWS 与 GCP 产品对应关系")
rows = [
    ("全球骨干", "AWS Global Network（默认全程自有骨干）", "Network Service Tiers（Premium/Standard 两档可选）"),
    ("VPC/内网", "VPC（区域级）· Peering · Transit Gateway · PrivateLink", "VPC（全局级）· Peering · NCC · Private Service Connect"),
    ("混合云互联", "Direct Connect · Site-to-Site VPN", "Cloud Interconnect · Cloud VPN（HA/Classic）"),
    ("负载均衡", "ELB（ALB/NLB/GWLB/CLB，区域级）", "Cloud Load Balancing（全局单 anycast IP）"),
    ("加速", "Global Accelerator（anycast 静态IP）", "无同名产品（Premium Tier + 全局 anycast LB）"),
    ("CDN", "CloudFront（一体，static+dynamic）", "Cloud CDN + Media CDN（拆分）"),
    ("DNS", "Route 53（8 种路由策略 + 域名注册）", "Cloud DNS（3 类路由策略）"),
    ("安全/DDoS", "SG+NACL · WAF · Shield Std/Advanced", "VPC Firewall · Cloud Armor（含 Adaptive）"),
]
two_col_table(s, rows, top=Inches(1.7), row_h=Inches(0.6), fs=11)
source_note(s, "详见后续各维度分页；每页脚注标注官方文档来源。")

# ============ Slide 3: 骨干网分隔 ============
sep_slide("维度 1 · 全球骨干网 & 网络分层", "cold potato vs hot potato：GCP 把网络分层做成可选产品，AWS 默认全程走自有骨干")

# ============ Slide 4: 骨干网 ============
s = add_slide(); bg(s, WHITE)
title_bar(s, "1. 全球骨干网 & 网络分层", "GCP 两档可选 vs AWS 默认单档自有骨干")
rows = [
    ("网络分层", "无用户可选分层（官方无该产品）", "Network Service Tiers：Premium / Standard 两档"),
    ("默认路由", "跨区流量始终走 AWS 骨干，从不经公网", "Premium=cold potato（尽早入网/晚出网）\nStandard=hot potato（尽早交 ISP）"),
    ("物理层加密", "所有数据中心间流量物理层自动加密", "Premium 骨干保护到最后一英里"),
    ("骨干 SLA", "无单独骨干 SLA（未取到，不臆造）", "Premium 99.99% / Standard 99.9%"),
    ("骨干规模", "近 2000 万公里光缆 · 39 Region/123 AZ", "官方公开光缆里程未取到（不确定）"),
    ("Standard 限制", "—", "仅区域IP/区域LB/Cloud NAT；不支持全局IP、CDN、Cloud VPN网关"),
]
two_col_table(s, rows, top=Inches(1.6), row_h=Inches(0.64), fs=11)
text(s, Inches(0.55), Inches(6.35), Inches(12.4), Inches(0.6),
     [("核心：GCP 把「全球骨干 vs ISP 出网」做成可计费选择；AWS 无等价分层，相当于只给 Premium 那档、不暴露降级省钱选项。", 11.5, MAIN_BLUE, True)])
source_note(s, "来源：cloud.google.com/network-tiers/docs/overview · docs.aws.amazon.com/vpc/latest/peering · aws.amazon.com/about-aws/global-infrastructure/")

# ============ Slide 5: VPC/内网 ============
s = add_slide(); bg(s, WHITE)
title_bar(s, "2. VPC / 内网架构 & 互联", "最根本差异：AWS VPC 区域级 vs GCP VPC 全局级")
rows = [
    ("VPC 作用域", "区域级（子网锁单 AZ）", "全局级（一个 VPC 子网可横跨多 region）"),
    ("Peering", "支持跨区走骨干加密；不传递路由", "延迟/吞吐同 VPC 内；官方明确不传递路由"),
    ("中心互联", "Transit Gateway（route table 显式路由\n+ Encryption control 强制加密）", "Network Connectivity Center（全局 hub\n+ 原生 Cross-Cloud 连 AWS/OCI）"),
    ("私有服务", "PrivateLink（消费者→服务）", "Private Service Connect（多出 Interfaces 反向双向）"),
    ("MTU", "TGW: VPC/DX 间 8500，VPN 1500", "Interconnect VLAN 可 8896 jumbo"),
]
two_col_table(s, rows, top=Inches(1.7), row_h=Inches(0.78), fs=11)
text(s, Inches(0.55), Inches(6.35), Inches(12.4), Inches(0.6),
     [("核心：模型层最大差异 = VPC 区域级 vs 全局级；NCC 多云 spoke + PSC 反向双向连接是 GCP 亮点。", 12, MAIN_BLUE, True)])
source_note(s, "来源：docs.cloud.google.com/vpc/docs/{vpc,vpc-peering,private-service-connect} · network-connectivity · docs.aws.amazon.com/vpc/latest/{userguide,peering,tgw,privatelink}")

# ============ Slide 6: 混合云互联 ============
s = add_slide(); bg(s, WHITE)
title_bar(s, "3. 混合云互联（专线 & VPN）", "Direct Connect vs Cloud Interconnect；S2S VPN vs Cloud VPN")
rows = [
    ("专线形态", "Dedicated / Hosted(Partner)", "Dedicated / Partner"),
    ("专线聚合", "最多 2×400G 或 4×<100G（LAG）", "Dedicated 到 8×400G=3200G；Partner 800G"),
    ("专线加密", "支持 MACsec（仅 dedicated+LAG，无额外费用）", "支持 MACsec（Dedicated/Partner 均可）"),
    ("专线 SLA", "多站冗余 99.99% / 单连接 ~95%", "生产级(冗余) 99.99%；单连接 No SLA"),
    ("VPN 结构", "单产品，每连接 2 隧道 HA", "HA VPN(99.99%) / Classic VPN(99.9%) 两代"),
    ("VPN 带宽/IPv6", "标准 1.25G，大带宽 5Gbps/隧道", "HA 支持双栈/IPv6；Classic 不支持 IPv6"),
]
two_col_table(s, rows, top=Inches(1.6), row_h=Inches(0.64), fs=11)
text(s, Inches(0.55), Inches(6.35), Inches(12.4), Inches(0.6),
     [("核心：专线结构对应，GCP 明确 3200G 聚合 + Partner L2/L3；VPN 上 GCP 用 HA/Classic 两代分明，AWS 单产品双隧道 + 5Gbps 大带宽隧道。", 11.5, MAIN_BLUE, True)])
source_note(s, "来源：docs.cloud.google.com/network-connectivity/docs/{interconnect,vpn} · docs.aws.amazon.com/directconnect（含 MACsec/lags/sla）· /vpn")

# ============ Slide 7: 外网入口分隔 ============
sep_slide("维度 4-5 · 外网入口 & 加速", "GCP 单 anycast IP 内建全球化 vs AWS 区域级 ELB + Global Accelerator 补全球")

# ============ Slide 8: 负载均衡 ============
s = add_slide(); bg(s, WHITE)
title_bar(s, "4. 负载均衡 & 外网入口", "GCP 全局单 anycast IP vs AWS 区域级 ELB")
rows = [
    ("L7", "ALB（HTTP/HTTPS/gRPC）", "Application LB（Global/Regional/Cross-region）"),
    ("L4 代理", "NLB（TCP/UDP/TLS，静态/弹性IP）", "Proxy Network LB（可选 SSL offload）"),
    ("L4 直通", "NLB 保留源IP", "Passthrough Network LB（UDP/ESP/GRE/ICMP）"),
    ("网关型", "GWLB（L3网关+L4）", "官方未列等价产品（不确定）"),
    ("全球化方式", "区域级，需叠加 Global Accelerator/Route53", "单个 anycast IP 全球分发（官方原文）"),
]
two_col_table(s, rows, top=Inches(1.7), row_h=Inches(0.72), fs=11)
text(s, Inches(0.55), Inches(6.15), Inches(12.4), Inches(0.6),
     [("核心：GCP 全局 LB = 单 anycast IP 内建全球分发（closest backend）；AWS ELB 区域级，全球入口需叠加 GA/Route 53。", 12, MAIN_BLUE, True)])
source_note(s, "来源：aws.amazon.com/elasticloadbalancing/features/ · docs.cloud.google.com/load-balancing/docs/choosing-load-balancer · cloud.google.com/load-balancing")

# ============ Slide 9: 加速 ============
s = add_slide(); bg(s, WHITE)
title_bar(s, "5. 加速", "AWS 有独立 Global Accelerator；GCP 无同名产品")
rows = [
    ("加速产品", "Global Accelerator（独立产品）", "无同名产品（Premium Tier + 全局 anycast LB 达成等价）"),
    ("IP", "2 个静态 anycast IPv4（独立 network zone）", "Premium Tier 全局 anycast IP"),
    ("路径", "全程走 AWS 骨干", "Premium Tier 走 Google 骨干"),
    ("协议/缓存", "TCP/UDP，不缓存（边缘丢 TCP 分片）", "依 LB 类型，不缓存（缓存靠 Cloud CDN）"),
    ("传输加速", "S3 Transfer Acceleration（走 CloudFront 边缘）", "无对标产品（仅 GCS 客户端并行/断点续传）"),
]
two_col_table(s, rows, top=Inches(1.7), row_h=Inches(0.72), fs=11)
text(s, Inches(0.55), Inches(6.15), Inches(12.4), Inches(0.6),
     [("核心：AWS 官方称 Global Accelerator 提速最高 60%；GCP 无独立加速器产品，靠 Premium Tier 功能等价（据实措辞）。", 12, MAIN_BLUE, True)])
source_note(s, "来源：docs.aws.amazon.com/global-accelerator/.../introduction-how-it-works.html · aws.amazon.com/global-accelerator/faqs/ · cloud.google.com/storage/docs/uploads")

# ============ Slide 10: CDN/DNS 分隔 ============
sep_slide("维度 6-7 · CDN & DNS", "CloudFront 一体通吃 vs Cloud CDN+Media CDN 拆分；Route 53 路由策略更丰富")

# ============ Slide 11: CDN ============
s = add_slide(); bg(s, WHITE)
title_bar(s, "6. CDN", "CloudFront 一体 vs GCP Cloud CDN + Media CDN 拆分")
rows = [
    ("定位", "CloudFront：Web 静态+动态通用", "Cloud CDN（挂 ALB）+ Media CDN（流媒体）"),
    ("独立性", "独立，源站 S3/HTTP", "Cloud CDN 必须挂 ALB，不能独立"),
    ("边缘节点", "750+ POPs + 15 Regional edge caches", "Cloud CDN 100+ locations（复用 Google PoP）"),
    ("定价", "Invalidation $0.005/条；AWS源→边缘免费\n(每GB分区单价官方表JS渲染未取到)", "cache egress $0.02–0.20/GiB\nMedia CDN 官方定价页 404 未取到"),
    ("YouTube 基础设施", "—", "官方无背书，只说 Google global edge-caching"),
]
two_col_table(s, rows, top=Inches(1.7), row_h=Inches(0.78), fs=10.5)
text(s, Inches(0.55), Inches(6.35), Inches(12.4), Inches(0.6),
     [("核心：CloudFront 单一产品通吃 static+dynamic；GCP 拆成 Cloud CDN（依附 ALB）+ Media CDN（流媒体）。", 12, MAIN_BLUE, True)])
source_note(s, "来源：docs.aws.amazon.com/AmazonCloudFront · aws.amazon.com/cloudfront/pricing · cloud.google.com/cdn/{docs/overview,docs/locations,pricing} · /media-cdn/docs/overview")

# ============ Slide 12: DNS ============
s = add_slide(); bg(s, WHITE)
title_bar(s, "7. DNS", "Route 53 路由策略更丰富（8 vs 3）+ 自带域名注册")
rows = [
    ("路由策略", "8 种：Simple/Failover/Geo/Geoproximity\n/Latency/IP-based/Multivalue/Weighted", "3 类：WRR / Geolocation(geofence) / Failover"),
    ("域名注册", "支持（三大功能之一）", "官方未提（不确定）"),
    ("DNSSEC", "支持 signing；KSK 需自管 KMS key", "支持；DNSKEY 创建+轮换+签名全自动托管"),
    ("SLA", "100% 承诺；最严重档赔 100% credit", "SLO 100%；credit 封顶 50%"),
    ("私有 DNS", "所有策略可用于 private zone", "部分 zone 类型不支持路由策略"),
]
two_col_table(s, rows, top=Inches(1.7), row_h=Inches(0.78), fs=10.5)
text(s, Inches(0.55), Inches(6.35), Inches(12.4), Inches(0.6),
     [("核心：Route 53 路由策略更细（Latency/Geoproximity/IP-based/Multivalue）+ 自带域名注册；DNSSEC Cloud DNS 更全托管。", 12, MAIN_BLUE, True)])
source_note(s, "来源：docs.aws.amazon.com/Route53/.../routing-policy.html + dnssec + aws.amazon.com/route53/sla · docs.cloud.google.com/dns/docs/{routing-policies,dnssec} · /dns/sla")

# ============ Slide 13: 安全分隔 ============
sep_slide("维度 8 · 网络安全 · WAF · DDoS", "防火墙深水区 · Cloud Armor vs Shield · DDoS 防护对比")

# ============ Slide 14: 防火墙深水区 ============
s = add_slide(); bg(s, WHITE)
title_bar(s, "8a. 防火墙层（深水区）", "SG（有状态）+ NACL（无状态） vs GCP 单一有状态防火墙")
rows = [
    ("有状态性", "SG 有状态 / NACL 无状态", "天生有状态（5-tuple，10 分钟活跃窗）"),
    ("作用层级", "SG 实例/ENI 级；NACL 子网级", "网络级定义，按 target(tag/SA) 应用到 VM"),
    ("动作/优先级", "SG 仅 allow 无优先级；\nNACL allow+deny 编号 1–32766 升序命中即停", "priority 0–65535（0最高）；\n同级冲突 deny 胜出；策略层 0–21亿"),
    ("默认姿态", "新建 SG 入站全拒/出站全放", "隐含 deny 全入站 + allow 全出站（65535 不可删）"),
    ("微隔离", "引用另一个 SG（sg-id 作源）", "network tag / service account 作 target 与 source"),
    ("关键配额", "规则×SG数 ≤1000；每ENI ≤16 SG", "连接跟踪封顶 1,040,000"),
]
two_col_table(s, rows, top=Inches(1.6), row_h=Inches(0.66), fs=10.5)
text(s, Inches(0.55), Inches(6.5), Inches(12.4), Inches(0.5),
     [("核心：AWS 两层（SG有状态实例级 + NACL无状态子网级）；GCP 只有一种、天生有状态 + 层级/全局/区域策略多级管线。", 11.5, MAIN_BLUE, True)])
source_note(s, "来源：docs.aws.amazon.com/vpc/latest/userguide/{vpc-security-groups,vpc-network-acls,infrastructure-security,amazon-vpc-limits} · cloud.google.com/firewall/docs")

# ============ Slide 15: WAF + NGFW ============
s = add_slide(); bg(s, WHITE)
title_bar(s, "8b. WAF & 下一代防火墙", "AWS WAF vs Cloud Armor；Network Firewall vs Cloud NGFW")
rows = [
    ("WAF 挂载点", "9 类资源（CloudFront/ALB/API GW 等）", "Cloud Armor（挂外部 ALB）"),
    ("WAF 规则", "托管规则集 + 原生 CAPTCHA/Challenge", "OWASP CRS 4.22 + Google 威胁情报"),
    ("NGFW 引擎", "Network Firewall：开源 Suricata IPS/IDS", "Cloud NGFW：L3/L4/L7 + IDS/IPS + URL filtering"),
    ("深度检测", "deep packet inspection + 按协议过滤 HTTPS", "分 Essentials 免费 / Standard 收费"),
    ("TLS 检测", "TLS inspection 解密→检测→再加密", "Standard 层 apply_security_profile_group"),
]
two_col_table(s, rows, top=Inches(1.7), row_h=Inches(0.72), fs=11)
text(s, Inches(0.55), Inches(6.15), Inches(12.4), Inches(0.6),
     [("核心：AWS WAF 挂载点更广 + 原生 CAPTCHA；Network Firewall 用 Suricata + TLS 解密检测；GCP Cloud NGFW 分免费/收费两档。", 11.5, MAIN_BLUE, True)])
source_note(s, "来源：docs.aws.amazon.com/waf · /network-firewall（Suricata/TLS inspection）· cloud.google.com/armor · cloud.google.com/firewall（Cloud NGFW）")

# ============ Slide 16: DDoS ============
s = add_slide(); bg(s, WHITE)
title_bar(s, "8c. DDoS 防护（重点）", "Shield Standard/Advanced vs Cloud Armor Standard/Enterprise")
rows = [
    ("免费层", "Shield Standard（默认，L3/L4 自动缓解）", "Cloud Armor Standard（L3/L4 自动缓解）"),
    ("免费层能力", "都不做签名匹配", "都不做签名匹配"),
    ("付费层", "Shield Advanced", "Cloud Armor Enterprise"),
    ("付费核心", "24/7 SRT 人肉响应 + 费用保护\n+ 500亿 WAF 请求/月额度 + L7 DDoS 托管规则", "ML 驱动 Adaptive Protection\n（Advanced network 训练 24h / Adaptive ≥1h）"),
    ("DDoS SLA/规模", "无 DDoS 专属 SLA、无官方可缓解规模数字", "无 DDoS 专属 SLA、无官方可缓解规模数字"),
]
two_col_table(s, rows, top=Inches(1.7), row_h=Inches(0.75), fs=10.5)
text(s, Inches(0.55), Inches(6.3), Inches(12.4), Inches(0.6),
     [("核心：Shield Advanced 强在 24/7 SRT + 费用保护 + WAF 请求额度；Cloud Armor Enterprise 强在 ML Adaptive Protection。两家均无 DDoS 专属 SLA。", 11, MAIN_BLUE, True)])
source_note(s, "来源：aws.amazon.com/shield（Standard/Advanced）· cloud.google.com/armor（Standard/Enterprise/Adaptive Protection）")

# ============ Slide 17: SLA 汇总 ============
s = add_slide(); bg(s, WHITE)
title_bar(s, "9. 稳定性 / SLA 汇总", "主要网络产品官方 SLA（均官方核实）")
rows = [
    ("负载均衡", "ELB Multi-AZ 99.99%", "Cloud LB Premium 99.99%"),
    ("专线", "Direct Connect 冗余 99.99% / 单连接 ~95%", "Cloud Interconnect 冗余 99.99% / 单连接 No SLA"),
    ("VPN", "S2S VPN（未取到专属 SLA 数字）", "HA VPN 99.99% / Classic VPN 99.9%"),
    ("DNS", "Route 53 100%（最严重赔 100%）", "Cloud DNS 100% SLO（封顶 50%）"),
    ("CDN", "CloudFront 99.9%", "Cloud CDN 独立 SLA 未确认（页 404）"),
    ("加速", "Global Accelerator 99.99% / 单AZ 99.5%", "—"),
]
two_col_table(s, rows, top=Inches(1.6), row_h=Inches(0.62), fs=11)
text(s, Inches(0.55), Inches(6.35), Inches(12.4), Inches(0.6),
     [("核心：主要产品 SLA 两家基本对齐（LB/专线 99.99%，DNS 100%）；GCP 专线单连接无 SLA 是选型注意点。", 11.5, MAIN_BLUE, True)])
source_note(s, "来源：各产品官方 SLA 页（aws.amazon.com/xxx/sla · cloud.google.com/xxx/sla），2026-08 核实。")

# ============ Slide 18: 总结 ============
s = add_slide(); bg(s, WHITE)
title_bar(s, "总结：关键差异一览", "选型时优先看这几点")
items = [
    ("全球骨干", "GCP 分层可选(Premium cold-potato/Standard hot-potato)；AWS 默认全程自有骨干无分层", AWS_ORANGE),
    ("VPC 模型", "AWS VPC 区域级 vs GCP VPC 全局级 —— 最根本的架构差异", GCP_BLUE),
    ("全球入口", "GCP 单 anycast IP 内建全球化 vs AWS 区域级 ELB + Global Accelerator", AWS_ORANGE),
    ("加速/CDN", "AWS 有独立 Global Accelerator + 一体 CloudFront；GCP 靠 Premium Tier + Cloud CDN/Media CDN 拆分", GREEN),
    ("DNS", "Route 53 路由策略更丰富(8 vs 3)+ 自带域名注册；两家 SLA 均 100% SLO", GCP_BLUE),
    ("防火墙", "AWS SG(有状态)+NACL(无状态)双层；GCP 单一有状态 + 层级策略 + tag/SA 微隔离", AWS_ORANGE),
    ("DDoS", "Shield Advanced=24/7 SRT+费用保护；Cloud Armor Enterprise=ML Adaptive；均无专属 SLA", GREEN),
]
y = Inches(1.5)
for title_t, body, col in items:
    box(s, Inches(0.55), y, Inches(0.15), Inches(0.72), fill=col)
    text(s, Inches(0.85), y, Inches(11.9), Inches(0.36), [(title_t, 14, col, True)])
    text(s, Inches(0.85), y + Inches(0.33), Inches(11.9), Inches(0.42), [(body, 11, DARK, False)])
    y += Inches(0.78)
source_note(s, "所有数据基于 AWS/GCP 官方文档核实；标「不确定/未取到」处为官方页本身未提供或无法抓取，绝不臆造。")

prs.save("AWS_vs_GCP_Network.pptx")
print("saved AWS_vs_GCP_Network.pptx", len(prs.slides.__iter__.__self__._sldIdLst), "slides")
