#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
ONTAP 存储层级关系 PPT (v2: 去掉 Plex, 加 FabricPool)
Disk -> RAID Group -> Aggregate (Local Tier) -> FlexVol Volume
       + FabricPool (高性能本地层 + 云对象存储容量层 自动分层)
基于 NetApp 官方文档 docs.netapp.com
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

NETAPP_BLUE = RGBColor(0x00, 0x67, 0xC5)
DARK = RGBColor(0x1A, 0x2B, 0x4A)
GRAY = RGBColor(0x5A, 0x6B, 0x82)
LIGHT_BG = RGBColor(0xF2, 0xF6, 0xFB)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
ORANGE = RGBColor(0xF5, 0x82, 0x20)
GREEN = RGBColor(0x2E, 0x9E, 0x5B)
ACCENT = RGBColor(0x0B, 0x5C, 0xAB)
CLOUD = RGBColor(0x6A, 0x5A, 0xCD)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
blank = prs.slide_layouts[6]


def add_slide():
    return prs.slides.add_slide(blank)


def bg(slide, color=WHITE):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    s.fill.solid(); s.fill.fore_color.rgb = color
    s.line.fill.background(); s.shadow.inherit = False
    slide.shapes._spTree.remove(s._element)
    slide.shapes._spTree.insert(2, s._element)
    return s


def box(slide, l, t, w, h, fill=None, line=None, line_w=1.0, shape=MSO_SHAPE.RECTANGLE):
    s = slide.shapes.add_shape(shape, l, t, w, h)
    if fill is None:
        s.fill.background()
    else:
        s.fill.solid(); s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line; s.line.width = Pt(line_w)
    s.shadow.inherit = False
    return s


def text(slide, l, t, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, wrap=True):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = Pt(2); tf.margin_right = Pt(2)
    tf.margin_top = Pt(1); tf.margin_bottom = Pt(1)
    if isinstance(runs, str):
        runs = [(runs, 18, DARK, False)]
    first = True
    for item in runs:
        txt, sz, col, bold = item[0], item[1], item[2], item[3]
        sp = item[4] if len(item) > 4 else None
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = align
        if sp is not None:
            p.space_before = Pt(sp)
        r = p.add_run(); r.text = txt
        r.font.size = Pt(sz); r.font.bold = bold
        r.font.color.rgb = col
        r.font.name = "Microsoft YaHei"
    return tb


def title_bar(slide, title, sub=None):
    box(slide, 0, 0, SW, Inches(1.05), fill=NETAPP_BLUE)
    box(slide, 0, Inches(1.05), SW, Pt(4), fill=ORANGE)
    text(slide, Inches(0.55), Inches(0.12), Inches(11), Inches(0.8),
         [(title, 28, WHITE, True)], anchor=MSO_ANCHOR.MIDDLE)
    if sub:
        text(slide, Inches(0.6), Inches(1.2), Inches(12), Inches(0.4),
             [(sub, 14, GRAY, False)])


# ============ Slide 1: 封面 ============
s = add_slide(); bg(s, DARK)
box(s, 0, Inches(2.9), SW, Pt(3), fill=ORANGE)
box(s, Inches(0.0), Inches(2.95), SW, Inches(0.02), fill=NETAPP_BLUE)
text(s, Inches(0.9), Inches(1.7), Inches(11.5), Inches(1.2),
     [("NetApp ONTAP 存储架构", 44, WHITE, True)])
text(s, Inches(0.9), Inches(3.15), Inches(11.5), Inches(1.0),
     [("Disk · RAID Group · Aggregate · Volume · FabricPool 之间的关系", 22, RGBColor(0x9F,0xC4,0xEA), False)])
text(s, Inches(0.9), Inches(6.5), Inches(11.5), Inches(0.5),
     [("资料来源：NetApp 官方文档 docs.netapp.com  |  2026-06", 13, GRAY, False)])
box(s, Inches(0.9), Inches(0.7), Inches(2.2), Inches(0.55), fill=NETAPP_BLUE)
text(s, Inches(0.9), Inches(0.7), Inches(2.2), Inches(0.55),
     [("NetApp ONTAP", 16, WHITE, True)], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


# ============ Slide 2: 总览层级图 (无 Plex) ============
s = add_slide(); bg(s, WHITE)
title_bar(s, "整体层级：从物理到逻辑", "物理磁盘逐层聚合，最终对外提供逻辑卷")

layers = [
    ("Disk / Drive（物理磁盘）", "SSD / HDD / 分区盘，最底层物理存储单元", RGBColor(0x6B,0x7C,0x93)),
    ("RAID Group（RAID 组）", "若干数据盘 + 校验盘，提供数据冗余保护", RGBColor(0x2E,0x80,0xC0)),
    ("Aggregate / Local Tier（聚合 / 本地层）", "ONTAP 核心物理存储池，由一个或多个 RAID 组构成", NETAPP_BLUE),
    ("FlexVol Volume（卷）", "驻留在 Aggregate 之上的逻辑卷，承载文件系统/LUN", GREEN),
]
top = Inches(1.55)
hgt = Inches(1.08)
gap = Inches(0.18)
fullw = Inches(8.4)
for i, (name, desc, col) in enumerate(layers):
    y = top + i * (hgt + gap)
    box(s, Inches(0.7), y, fullw, hgt, fill=col, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    text(s, Inches(0.9), y, fullw-Inches(0.4), hgt,
         [(name, 19, WHITE, True), (desc, 12, RGBColor(0xEA,0xF2,0xFB), False, 3)],
         anchor=MSO_ANCHOR.MIDDLE)
    if i < len(layers)-1:
        text(s, Inches(0.7)+fullw/2-Inches(0.3), y+hgt-Inches(0.02), Inches(0.6), gap+Inches(0.15),
             [("▼", 16, GRAY, True)], align=PP_ALIGN.CENTER)

rx = Inches(9.5)
box(s, rx, Inches(1.55), Inches(3.2), Inches(5.0), fill=LIGHT_BG, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
text(s, rx+Inches(0.2), Inches(1.75), Inches(2.85), Inches(4.7),
     [("聚合方向", 16, NETAPP_BLUE, True),
      ("自下而上：物理 → 逻辑", 12, DARK, False, 6),
      ("• 磁盘组成 RAID 组", 12, DARK, False, 10),
      ("• RAID 组组成 Aggregate", 12, DARK, False, 6),
      ("• Volume 建在 Aggregate 上", 12, DARK, False, 6),
      ("一个 Aggregate 可承载", 12, GRAY, False, 16),
      ("多个 FlexVol", 13, ORANGE, True, 2),
      ("一个 Volume 只属于", 12, GRAY, False, 12),
      ("一个 Aggregate", 13, ORANGE, True, 2),
      ("FabricPool 可为聚合", 12, GRAY, False, 14),
      ("挂接云对象存储容量层", 13, CLOUD, True, 2)])


# ============ Slide 3: Disk & RAID Group ============
s = add_slide(); bg(s, WHITE)
title_bar(s, "Disk 与 RAID Group", "磁盘是物理单元，RAID 组提供冗余")

lx = Inches(0.7)
box(s, lx, Inches(1.45), Inches(5.7), Inches(5.4), fill=LIGHT_BG, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
text(s, lx+Inches(0.3), Inches(1.6), Inches(5.1), Inches(0.5),
     [("Disk / Drive（磁盘）", 20, NETAPP_BLUE, True)])
text(s, lx+Inches(0.3), Inches(2.25), Inches(5.1), Inches(4.4),
     [("• ONTAP 管理的最底层物理存储单元", 14, DARK, False),
      ("• 类型：SSD、SAS/HDD、容量盘", 14, DARK, False, 8),
      ("• 可被「分区」共享给多个 RAID 组", 14, DARK, False, 8),
      ("   （Position 显示 shared）", 12, GRAY, False, 2),
      ("• 磁盘角色（Position）：", 14, DARK, False, 12),
      ("    data   —  数据盘", 13, GREEN, True, 4),
      ("    parity —  行校验盘", 13, ORANGE, True, 2),
      ("    dparity — 对角校验盘（RAID-DP）", 13, ORANGE, True, 2),
      ("• 未分配磁盘 = spare（热备盘）", 14, DARK, False, 12),
      ("• 命令：storage disk show", 12, ACCENT, False, 12)])

rx = Inches(6.8)
box(s, rx, Inches(1.45), Inches(5.85), Inches(5.4), fill=WHITE, line=NETAPP_BLUE, line_w=1.5, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
text(s, rx+Inches(0.3), Inches(1.6), Inches(5.3), Inches(0.5),
     [("RAID Group（RAID 组，以 RAID-DP 为例）", 18, NETAPP_BLUE, True)])
disk_specs = [("D", GREEN)]*4 + [("P", ORANGE), ("dP", ORANGE)]
dx0 = rx + Inches(0.35); dy0 = Inches(2.5)
dw = Inches(0.78); dh = Inches(1.0); dgap = Inches(0.12)
for i, (lab, col) in enumerate(disk_specs):
    x = dx0 + i*(dw+dgap)
    box(s, x, dy0, dw, dh, fill=col, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    text(s, x, dy0, dw, dh, [(lab, 18, WHITE, True)], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
text(s, rx+Inches(0.35), dy0+dh+Inches(0.1), Inches(5.3), Inches(0.4),
     [("D = 数据盘    P = 校验盘    dP = 双校验盘", 12, GRAY, False)])
text(s, rx+Inches(0.3), Inches(4.3), Inches(5.3), Inches(2.4),
     [("• RAID 组 = 一组数据盘 + 校验盘", 14, DARK, False),
      ("• 常见类型：RAID-DP（双校验，默认）、", 14, DARK, False, 8),
      ("   RAID-TEC（三校验）、RAID4", 14, DARK, False, 2),
      ("• 校验盘提供磁盘故障时的数据重建能力", 14, DARK, False, 8),
      ("• RAID-DP 可同时容忍 2 块盘故障", 14, GREEN, True, 8),
      ("• 多个 RAID 组组成一个 Aggregate", 14, DARK, False, 8),
      ("• 命令：storage aggregate show-status", 12, ACCENT, False, 10)])


# ============ Slide 4: Aggregate ============
s = add_slide(); bg(s, WHITE)
title_bar(s, "Aggregate（聚合 / 本地层）", "ONTAP 核心物理存储池，由多个 RAID 组构成")

lx = Inches(0.7)
box(s, lx, Inches(1.45), Inches(5.7), Inches(5.4), fill=LIGHT_BG, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
text(s, lx+Inches(0.3), Inches(1.6), Inches(5.1), Inches(0.5),
     [("Aggregate 关键概念", 20, NETAPP_BLUE, True)])
text(s, lx+Inches(0.3), Inches(2.25), Inches(5.1), Inches(4.4),
     [("• ONTAP 核心物理存储池", 14, DARK, False),
      ("• 由一个或多个 RAID 组聚合而成", 14, DARK, False, 8),
      ("• 9.7 之前 System Manager 称「本地层」", 13, GRAY, False, 8),
      ("   CLI 始终称 aggregate", 13, GRAY, False, 2),
      ("• 是创建 FlexVol 卷的基础容器", 14, DARK, False, 8),
      ("• 一个 Aggregate 可承载多个 FlexVol", 14, GREEN, True, 8),
      ("• 可按性能需求隔离工作负载：", 14, DARK, False, 8),
      ("   全 SSD 聚合 → 低延迟高性能", 12, GRAY, False, 2),
      ("   HDD/容量盘聚合 → 归档冷数据", 12, GRAY, False, 2),
      ("• 通过 FabricPool 可挂接云容量层", 14, CLOUD, True, 8),
      ("• 命令：storage aggregate show", 12, ACCENT, False, 8)])

rx = Inches(6.7)
box(s, rx, Inches(1.45), Inches(5.95), Inches(5.4), fill=WHITE, line=NETAPP_BLUE, line_w=1.5, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
text(s, rx+Inches(0.25), Inches(1.55), Inches(5.4), Inches(0.4),
     [("Aggregate 结构示意", 16, NETAPP_BLUE, True)])
box(s, rx+Inches(0.35), Inches(2.05), Inches(5.2), Inches(4.6), fill=RGBColor(0xE8,0xF1,0xFB), line=NETAPP_BLUE, line_w=1.5, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
text(s, rx+Inches(0.5), Inches(2.12), Inches(4.8), Inches(0.35),
     [("Aggregate (aggr1)", 14, NETAPP_BLUE, True)])
for j, ry in enumerate([Inches(2.6), Inches(4.55)]):
    box(s, rx+Inches(0.6), ry, Inches(4.7), Inches(1.75), fill=WHITE, line=ORANGE, line_w=1.0, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    text(s, rx+Inches(0.75), ry+Inches(0.08), Inches(4.0), Inches(0.3),
         [("RAID Group rg%d (RAID-DP)" % j, 12, ORANGE, True)])
    specs = [GREEN,GREEN,GREEN,ORANGE,ORANGE]
    sx0 = rx+Inches(0.85); sy = ry+Inches(0.6); sw=Inches(0.7); sh2=Inches(0.95); sg=Inches(0.13)
    for k,col in enumerate(specs):
        x = sx0 + k*(sw+sg)
        box(s, x, sy, sw, sh2, fill=col, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        text(s, x, sy, sw, sh2, [("disk", 10, WHITE, True)], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


# ============ Slide 5: Volume ============
s = add_slide(); bg(s, WHITE)
title_bar(s, "Volume（FlexVol 卷）", "驻留在 Aggregate 之上的逻辑存储容器")

lx = Inches(0.7)
box(s, lx, Inches(1.45), Inches(6.0), Inches(5.4), fill=LIGHT_BG, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
text(s, lx+Inches(0.3), Inches(1.6), Inches(5.4), Inches(0.5),
     [("FlexVol 卷的特性", 20, NETAPP_BLUE, True)])
text(s, lx+Inches(0.3), Inches(2.25), Inches(5.4), Inches(4.5),
     [("• 逻辑存储容器，承载文件系统 / LUN", 14, DARK, False),
      ("• 必须建立在某一个 Aggregate 之上", 14, DARK, False, 8),
      ("• 一个 Aggregate 可放多个 FlexVol", 14, GREEN, True, 8),
      ("• 一个 Volume 只属于一个 Aggregate", 14, ORANGE, True, 6),
      ("• 可在线扩容/缩容，与物理盘解耦", 14, DARK, False, 8),
      ("• 支持精简配置（Thin Provisioning）：", 14, DARK, False, 8),
      ("    创建时不预留空间，按需占用", 12, GRAY, False, 2),
      ("    可超额配置（overcommit）聚合", 12, GRAY, False, 2),
      ("• 通过 SVM（Vserver）对外提供 NFS/", 14, DARK, False, 8),
      ("   CIFS/iSCSI/FC 访问", 14, DARK, False, 2),
      ("• 命令：volume show / volume create", 12, ACCENT, False, 8)])

rx = Inches(7.0)
box(s, rx, Inches(1.45), Inches(5.65), Inches(5.4), fill=WHITE, line=NETAPP_BLUE, line_w=1.5, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
text(s, rx+Inches(0.25), Inches(1.55), Inches(5.1), Inches(0.4),
     [("一个 Aggregate 承载多个卷", 16, NETAPP_BLUE, True)])
for j,(nm,col) in enumerate([("vol_nfs", GREEN),("vol_cifs", RGBColor(0x2E,0x80,0xC0)),("vol_lun (SAN)", RGBColor(0x5A,0x5A,0xC0))]):
    y = Inches(2.15) + j*Inches(0.95)
    box(s, rx+Inches(0.4), y, Inches(4.85), Inches(0.78), fill=col, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    text(s, rx+Inches(0.4), y, Inches(4.85), Inches(0.78),
         [("FlexVol: %s" % nm, 14, WHITE, True)], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
text(s, rx+Inches(0.4), Inches(5.0), Inches(4.85), Inches(0.3),
     [("▲ 逻辑层  |  ▼ 物理层", 11, GRAY, True)], align=PP_ALIGN.CENTER)
box(s, rx+Inches(0.4), Inches(5.4), Inches(4.85), Inches(1.2), fill=NETAPP_BLUE, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
text(s, rx+Inches(0.4), Inches(5.4), Inches(4.85), Inches(1.2),
     [("Aggregate（聚合物理池）", 16, WHITE, True),
      ("RAID Group → Disk", 12, RGBColor(0xCF,0xE3,0xF7), False, 4)],
     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


# ============ Slide 6: FabricPool ============
s = add_slide(); bg(s, WHITE)
title_bar(s, "FabricPool（自动分层）", "热数据留本地 SSD，冷数据自动下沉到云对象存储")

# 左侧说明
lx = Inches(0.7)
box(s, lx, Inches(1.45), Inches(5.7), Inches(5.4), fill=LIGHT_BG, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
text(s, lx+Inches(0.3), Inches(1.6), Inches(5.1), Inches(0.5),
     [("FabricPool 是什么", 20, NETAPP_BLUE, True)])
text(s, lx+Inches(0.3), Inches(2.25), Inches(5.1), Inches(4.5),
     [("• ONTAP 自动分层特性，按访问频率在", 14, DARK, False),
      ("   高性能本地层与云容量层间移动数据", 14, DARK, False, 2),
      ("• 热数据（hot）留在本地 SSD 性能层", 14, GREEN, True, 8),
      ("• 冷数据（cold）自动下沉到云对象存储", 14, CLOUD, True, 6),
      ("• 工作在存储块级别，文件和 LUN 均支持", 14, DARK, False, 8),
      ("• 应用无感知，分层全自动，无需运维干预", 14, DARK, False, 8),
      ("• 云容量层可选：", 14, DARK, False, 8),
      ("   AWS S3 / Azure Blob / GCP / 阿里云等", 12, GRAY, False, 2),
      ("   或 NetApp StorageGRID（私有云）", 12, GRAY, False, 2),
      ("• 释放昂贵本地容量，构建混合云", 14, ORANGE, True, 8),
      ("• 命令：storage aggregate object-store", 12, ACCENT, False, 8)])

# 右侧分层图
rx = Inches(6.7)
box(s, rx, Inches(1.45), Inches(5.95), Inches(5.4), fill=WHITE, line=NETAPP_BLUE, line_w=1.5, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
text(s, rx+Inches(0.25), Inches(1.55), Inches(5.4), Inches(0.4),
     [("FabricPool 分层示意", 16, NETAPP_BLUE, True)])

# 性能层（本地 SSD aggregate）
box(s, rx+Inches(0.45), Inches(2.15), Inches(5.05), Inches(1.55), fill=RGBColor(0xE3,0xF3,0xE9), line=GREEN, line_w=1.5, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
text(s, rx+Inches(0.6), Inches(2.25), Inches(4.7), Inches(0.4),
     [("性能层 Performance Tier（本地 SSD Aggregate）", 13, GREEN, True)])
for k in range(4):
    x = rx+Inches(0.7) + k*Inches(1.12)
    box(s, x, Inches(2.75), Inches(0.95), Inches(0.8), fill=GREEN, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    text(s, x, Inches(2.75), Inches(0.95), Inches(0.8), [("热数据", 11, WHITE, True)], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

# 双向箭头
text(s, rx+Inches(0.45), Inches(3.78), Inches(5.05), Inches(0.5),
     [("⇅  自动分层（基于访问温度，块级别）", 13, ORANGE, True)], align=PP_ALIGN.CENTER)

# 容量层（云对象存储）
box(s, rx+Inches(0.45), Inches(4.35), Inches(5.05), Inches(2.0), fill=RGBColor(0xEC,0xE9,0xFA), line=CLOUD, line_w=1.5, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
text(s, rx+Inches(0.6), Inches(4.45), Inches(4.7), Inches(0.4),
     [("容量层 Capacity Tier（云对象存储 / StorageGRID）", 13, CLOUD, True)])
for k in range(4):
    x = rx+Inches(0.7) + k*Inches(1.12)
    box(s, x, Inches(4.95), Inches(0.95), Inches(0.8), fill=CLOUD, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    text(s, x, Inches(4.95), Inches(0.95), Inches(0.8), [("冷数据", 11, WHITE, True)], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
text(s, rx+Inches(0.6), Inches(5.85), Inches(4.7), Inches(0.4),
     [("S3 / Blob / GCP / 阿里云 OSS / StorageGRID", 11, GRAY, False)], align=PP_ALIGN.CENTER)


# ============ Slide 7: 总结表 + 命令 ============
s = add_slide(); bg(s, WHITE)
title_bar(s, "层级关系总结", "概念 · 组成 · 常用 CLI 命令")

rows = [
    ("层级", "定义", "由什么组成", "查看命令"),
    ("Disk", "物理磁盘 (SSD/HDD)", "—（最底层）", "storage disk show"),
    ("RAID Group", "冗余保护单元", "数据盘 + 校验盘", "storage aggregate\nshow-status"),
    ("Aggregate", "物理存储池/本地层", "1+ 个 RAID Group", "storage aggregate show"),
    ("Volume", "逻辑卷 (FlexVol)", "驻留于 1 个 Aggregate", "volume show"),
    ("FabricPool", "聚合的自动分层", "本地 SSD 层 + 云对象层", "storage aggregate\nobject-store show"),
]
tx = Inches(0.6); ty = Inches(1.45)
colw = [Inches(2.1), Inches(3.3), Inches(3.6), Inches(3.4)]
rh = Inches(0.82)
y = ty
for ri, row in enumerate(rows):
    x = tx
    header = ri == 0
    for ci, cell in enumerate(row):
        fill = NETAPP_BLUE if header else (LIGHT_BG if ri % 2 else WHITE)
        box(s, x, y, colw[ci], rh, fill=fill, line=RGBColor(0xD5,0xDE,0xEA), line_w=0.75)
        col = WHITE if header else (NETAPP_BLUE if ci==0 else DARK)
        sz = 14 if header else (13 if ci==0 else 12)
        bold = header or ci==0
        text(s, x+Inches(0.1), y, colw[ci]-Inches(0.2), rh,
             [(cell, sz, col, bold)], anchor=MSO_ANCHOR.MIDDLE,
             align=PP_ALIGN.CENTER if (header or ci>=2) else PP_ALIGN.LEFT)
        x += colw[ci]
    y += rh

text(s, Inches(0.6), Inches(6.5), Inches(12.1), Inches(0.7),
     [("一句话关系：磁盘 → 组成 RAID 组 → 组成 Aggregate → 之上创建 Volume → 经 SVM 对外提供访问；FabricPool 为聚合挂接云对象存储，自动冷热分层",
       13, ACCENT, True)])

prs.save("/home/ubuntu/.openclaw/workspace/ONTAP_存储架构_Disk_Aggregate_Volume.pptx")
print("saved")
