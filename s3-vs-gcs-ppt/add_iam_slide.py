#!/usr/bin/env python3
"""在现有 S3_vs_GCS.pptx 追加一页「⑫ 访问控制 & IAM 对比」，插到最后总结页之前，保持原样式。"""
import copy
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

import os
SRC = os.path.join(os.path.dirname(os.path.abspath(__file__)), 'S3_vs_GCS.pptx')

DARK   = RGBColor(0x23, 0x2F, 0x3E)  # AWS 深蓝标题栏
ORANGE = RGBColor(0xFF, 0x85, 0x00)  # 橙分隔线/要点
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
SUBTXT = RGBColor(0xC9, 0xD2, 0xDE)  # 副标题浅灰蓝
BODY   = RGBColor(0x1A, 0x2B, 0x4A)
PAGENO = RGBColor(0x5A, 0x6B, 0x82)
ROWALT = RGBColor(0xF2, 0xF6, 0xFB)
FONT   = 'Amazon Ember'
FONT_H = 'Amazon Ember Heavy'

prs = Presentation(SRC)
blank = prs.slide_layouts[6]  # Blank
slide = prs.slides.add_slide(blank)

def add_rect(l,t,w,h,color):
    from pptx.enum.shapes import MSO_SHAPE
    sp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(l),Inches(t),Inches(w),Inches(h))
    sp.fill.solid(); sp.fill.fore_color.rgb=color
    sp.line.fill.background()
    sp.shadow.inherit=False
    return sp

def add_text(l,t,w,h,runs,size,bold=False,color=BODY,font=FONT,align=PP_ALIGN.LEFT,anchor=MSO_ANCHOR.TOP):
    tb=slide.shapes.add_textbox(Inches(l),Inches(t),Inches(w),Inches(h))
    tf=tb.text_frame; tf.word_wrap=True
    tf.vertical_anchor=anchor
    if isinstance(runs,str): runs=[(runs,color,bold)]
    p=tf.paragraphs[0]; p.alignment=align
    for txt,c,b in runs:
        r=p.add_run(); r.text=txt
        r.font.name=font; r.font.size=Pt(size); r.font.bold=b; r.font.color.rgb=c
    return tb

# 标题栏
add_rect(0,0,13.33,1.05,DARK)
add_rect(0,1.05,13.33,0.06,ORANGE)
add_text(0.55,0.12,12.2,0.66,'⑫ 访问控制 & IAM 对比',28,True,WHITE,FONT_H)
add_text(0.55,0.72,12.2,0.30,'授权模型 · 前缀级授权 · 临时访问 —— 基于官方文档核实',14,False,SUBTXT,FONT)

# 表格
rows=[
 ('维度','AWS S3','GCP Cloud Storage'),
 ('授权模型','IAM Policy + Bucket Policy(基于资源) + ACL','IAM(含 IAM Conditions) + ACL(legacy)；无 bucket policy'),
 ('前缀级授权','Bucket Policy 直接写 Resource: bucket/prefix/*(原生直接)','IAM Conditions 的 resource.name 前缀条件(官方支持,机制不同)'),
 ('统一 vs 细粒度','Bucket owner enforced(关闭 ACL,只用 IAM)','Uniform bucket-level access vs Fine-grained(IAM+ACL)'),
 ('临时访问','Presigned URL','Signed URL / Signed Policy Document'),
 ('防止公开','Block Public Access(账号/桶级)','Public Access Prevention(组织策略可强制)'),
 ('组织级治理','SCP + IAM(账号边界)','IAM 层级继承(Org→Folder→Project→Bucket)'),
]
nr=len(rows); nc=3
gt=slide.shapes.add_table(nr,nc,Inches(0.55),Inches(1.35),Inches(12.25),Inches(3.9)).table
gt.columns[0].width=Inches(2.5); gt.columns[1].width=Inches(4.6); gt.columns[2].width=Inches(5.15)
for ri,row in enumerate(rows):
    for ci,val in enumerate(row):
        cell=gt.cell(ri,ci)
        cell.margin_left=Inches(0.08); cell.margin_right=Inches(0.08)
        cell.margin_top=Inches(0.03); cell.margin_bottom=Inches(0.03)
        cell.vertical_anchor=MSO_ANCHOR.MIDDLE
        if ri==0:
            cell.fill.solid(); cell.fill.fore_color.rgb=DARK
        else:
            cell.fill.solid(); cell.fill.fore_color.rgb=(ROWALT if ri%2==0 else WHITE)
        tf=cell.text_frame; tf.word_wrap=True
        p=tf.paragraphs[0]
        r=p.add_run(); r.text=val
        r.font.name=FONT; r.font.size=Pt(11.5 if ri else 12.5)
        if ri==0:
            r.font.bold=True; r.font.color.rgb=WHITE
        elif ci==0:
            r.font.bold=True; r.font.color.rgb=DARK
        else:
            r.font.color.rgb=BODY

# 要点
notes=[
 '▸ 最大差异：S3 有「基于资源的 Bucket Policy」可直接按前缀(prefix/*)授权；GCS 无 bucket policy,前缀授权改用 IAM Conditions 的 resource.name 条件(官方支持,CEL 表达式)',
 '▸ 别把 GCS 说成「像 S3 那样按前缀配 policy」——机制不同：S3=桶上挂 JSON policy；GCS=IAM 条件式角色绑定',
 '▸ 两家都在推「只用 IAM、关闭对象级 ACL」的统一模型(S3 Bucket owner enforced / GCS Uniform bucket-level access)',
]
tb=slide.shapes.add_textbox(Inches(0.55),Inches(5.45),Inches(12.25),Inches(1.55))
tf=tb.text_frame; tf.word_wrap=True
for i,n in enumerate(notes):
    p=tf.paragraphs[0] if i==0 else tf.add_paragraph()
    r=p.add_run(); r.text=n
    r.font.name=FONT; r.font.size=Pt(11.5); r.font.bold=True; r.font.color.rgb=ORANGE
    p.space_after=Pt(4)

# 页码
add_text(12.35,7.03,0.75,0.35,'12',11,False,PAGENO,FONT,PP_ALIGN.RIGHT)

# 把新页移动到最后一页(总结页)之前
xml=prs.slides._sldIdLst
ids=list(xml)
new=ids[-1]           # 刚加的
xml.remove(new)
xml.insert(len(ids)-2, new)  # 插到原最后一页(总结)之前

prs.save(SRC)
print('done, slides:', len(prs.slides.__iter__.__self__._sldIdLst))
