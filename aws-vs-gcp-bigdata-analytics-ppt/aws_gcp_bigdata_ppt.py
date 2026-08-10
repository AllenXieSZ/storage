#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""AWS vs GCP 大数据/分析产品对比 PPT (NetApp风)"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# 配色
BLUE   = RGBColor(0x00,0x67,0xC5)
ORANGE = RGBColor(0xF5,0x82,0x20)
GREEN  = RGBColor(0x2E,0x9E,0x5B)
DARK   = RGBColor(0x1A,0x2B,0x4A)
LIGHT  = RGBColor(0xF2,0xF6,0xFB)
GREY   = RGBColor(0x5A,0x6B,0x82)
WHITE  = RGBColor(0xFF,0xFF,0xFF)
AWSCOL = RGBColor(0xFF,0x99,0x00)   # AWS橙
GCPCOL = RGBColor(0x42,0x85,0xF4)   # GCP蓝
FONT   = "Microsoft YaHei"

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
blank = prs.slide_layouts[6]

def _set(run, size, color, bold=False):
    run.font.size=Pt(size); run.font.color.rgb=color; run.font.bold=bold; run.font.name=FONT

def box(slide,l,t,w,h,fill=None,line=None,lw=1.0,round=False):
    from pptx.enum.shapes import MSO_SHAPE
    shp=slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if round else MSO_SHAPE.RECTANGLE,l,t,w,h)
    if fill is None: shp.fill.background()
    else: shp.fill.solid(); shp.fill.fore_color.rgb=fill
    if line is None: shp.line.fill.background()
    else: shp.line.color.rgb=line; shp.line.width=Pt(lw)
    shp.shadow.inherit=False
    return shp

def txt(slide,l,t,w,h,lines,align=PP_ALIGN.LEFT,anchor=MSO_ANCHOR.TOP,wrap=True):
    tb=slide.shapes.add_textbox(l,t,w,h); tf=tb.text_frame
    tf.word_wrap=wrap; tf.vertical_anchor=anchor
    tf.margin_left=Inches(0.05);tf.margin_right=Inches(0.05)
    tf.margin_top=Inches(0.02);tf.margin_bottom=Inches(0.02)
    for i,(t_,s,c,b) in enumerate(lines):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph()
        p.alignment=align; r=p.add_run(); r.text=t_; _set(r,s,c,b)
    return tb

def header(slide,title,sub):
    box(slide,0,0,SW,Inches(1.05),fill=BLUE)
    box(slide,0,Inches(1.05),SW,Pt(4),fill=ORANGE)
    txt(slide,Inches(0.5),Inches(0.12),Inches(12.3),Inches(0.55),[(title,26,WHITE,True)],anchor=MSO_ANCHOR.MIDDLE)
    txt(slide,Inches(0.5),Inches(0.66),Inches(12.3),Inches(0.32),[(sub,13,RGBColor(0xCF,0xE0,0xF5),False)],anchor=MSO_ANCHOR.MIDDLE)

# ---------- 封面 ----------
s=prs.slides.add_slide(blank)
box(s,0,0,SW,SH,fill=DARK)
box(s,0,0,Inches(0.28),SH,fill=ORANGE)
box(s,Inches(0.9),Inches(2.5),Inches(7.2),Pt(3),fill=ORANGE)
txt(s,Inches(0.9),Inches(1.5),Inches(11.5),Inches(1.0),[("AWS vs GCP",44,WHITE,True)])
txt(s,Inches(0.9),Inches(2.7),Inches(11.5),Inches(1.2),[("大数据与分析产品全景对比",34,WHITE,True)])
txt(s,Inches(0.9),Inches(3.9),Inches(11.5),Inches(0.9),[
    ("10 大类别 · 逐产品优缺点 · 湖仓/表格式(含 S3 Tables)",17,RGBColor(0x9F,0xB4,0xD4),False)])
txt(s,Inches(0.9),Inches(6.3),Inches(11.5),Inches(0.7),[
    ("基于 AWS / GCP 官方文档 · 2026-08 · 数据随产品更新，以官方为准",12,GREY,False)])

# ---------- 概览页 ----------
s=prs.slides.add_slide(blank)
header(s,"总览：两种架构哲学","AWS 多而全(purpose-built) vs GCP 少而精(serverless为主)")
rows=[
 ("架构哲学","每种负载一个专门服务，组合灵活、生态广","BigQuery 一个平台包揽多数分析，免运维"),
 ("计费主基调","按实例/节点时长 或 serverless按用量","serverless 为主(扫描量/slot)，无需管集群"),
 ("独有强项","搜索(OpenSearch)、托管Kafka(MSK)、精细选型","数仓(BigQuery)、批流一体(Dataflow/Beam)、语义层(Looker)"),
 ("开放格式","S3 Tables 原生 Iceberg、引擎中立","Iceberg/Delta/Hudi 内建，但引擎偏 BigQuery"),
]
y=Inches(1.35); ch=Inches(1.32)
# 表头
box(s,Inches(0.5),y,Inches(2.4),Inches(0.5),fill=DARK)
box(s,Inches(2.9),y,Inches(5.0),Inches(0.5),fill=AWSCOL)
box(s,Inches(7.9),y,Inches(4.9),Inches(0.5),fill=GCPCOL)
txt(s,Inches(0.5),y,Inches(2.4),Inches(0.5),[("维度",13,WHITE,True)],PP_ALIGN.CENTER,MSO_ANCHOR.MIDDLE)
txt(s,Inches(2.9),y,Inches(5.0),Inches(0.5),[("AWS",13,WHITE,True)],PP_ALIGN.CENTER,MSO_ANCHOR.MIDDLE)
txt(s,Inches(7.9),y,Inches(4.9),Inches(0.5),[("GCP",13,WHITE,True)],PP_ALIGN.CENTER,MSO_ANCHOR.MIDDLE)
y=y+Inches(0.5)
for i,(d,a,g) in enumerate(rows):
    bg=LIGHT if i%2==0 else WHITE
    box(s,Inches(0.5),y,Inches(2.4),ch,fill=bg,line=RGBColor(0xDD,0xE5,0xF0))
    box(s,Inches(2.9),y,Inches(5.0),ch,fill=bg,line=RGBColor(0xDD,0xE5,0xF0))
    box(s,Inches(7.9),y,Inches(4.9),ch,fill=bg,line=RGBColor(0xDD,0xE5,0xF0))
    txt(s,Inches(0.6),y,Inches(2.2),ch,[(d,12,BLUE,True)],PP_ALIGN.LEFT,MSO_ANCHOR.MIDDLE)
    txt(s,Inches(3.0),y,Inches(4.8),ch,[(a,11.5,DARK,False)],PP_ALIGN.LEFT,MSO_ANCHOR.MIDDLE)
    txt(s,Inches(8.0),y,Inches(4.7),ch,[(g,11.5,DARK,False)],PP_ALIGN.LEFT,MSO_ANCHOR.MIDDLE)
    y=y+ch

# ---------- 改名提示页 ----------
s=prs.slides.add_slide(blank)
header(s,"⚠️ 4 个近期改名(已核实官方文档)","做技术选型/沟通时用新名，避免踩坑")
renames=[
 ("Dataproc","→  Google Cloud Managed Service for Apache Spark","官方文档已重定向，旧名仍在用"),
 ("QuickSight","→  Amazon Quick Sight","隶属新品牌 Amazon Quick(AI套件)，原 Q 自然语言并入"),
 ("Cloud Composer","→  Managed Airflow (Gen3)","文档新称谓"),
 ("BigLake Iceberg tables","→  Apache Iceberg managed tables","BigQuery 官方新称谓"),
]
y=Inches(1.5)
for old,new,note in renames:
    box(s,Inches(0.7),y,Inches(11.9),Inches(1.15),fill=LIGHT,line=ORANGE,lw=1.5,round=True)
    txt(s,Inches(1.0),y+Inches(0.13),Inches(11.4),Inches(0.45),[(old+"  "+new,15,DARK,True)])
    txt(s,Inches(1.0),y+Inches(0.63),Inches(11.4),Inches(0.4),[(note,12,GREY,False)])
    y=y+Inches(1.35)

# ---------- 类别对比页模板 ----------
def cat_slide(num,cat,aws_name,gcp_name,aws_pos,gcp_pos,aws_pro,aws_con,gcp_pro,gcp_con,advice):
    s=prs.slides.add_slide(blank)
    header(s,f"{num}. {cat}",f"{aws_name}  vs  {gcp_name}")
    colw=Inches(6.0); gap=Inches(0.4); lx=Inches(0.45); rx=lx+colw+gap
    top=Inches(1.3)
    # AWS 卡
    box(s,lx,top,colw,Inches(0.5),fill=AWSCOL,round=False)
    txt(s,lx,top,colw,Inches(0.5),[("AWS · "+aws_name,14,WHITE,True)],PP_ALIGN.CENTER,MSO_ANCHOR.MIDDLE)
    box(s,lx,top+Inches(0.5),colw,Inches(4.85),fill=LIGHT,line=RGBColor(0xE0,0xC0,0x90))
    # GCP 卡
    box(s,rx,top,colw,Inches(0.5),fill=GCPCOL)
    txt(s,rx,top,colw,Inches(0.5),[("GCP · "+gcp_name,14,WHITE,True)],PP_ALIGN.CENTER,MSO_ANCHOR.MIDDLE)
    box(s,rx,top+Inches(0.5),colw,Inches(4.85),fill=LIGHT,line=RGBColor(0xB8,0xCE,0xF5))
    def fill_card(x,pos,pro,con):
        yy=top+Inches(0.62)
        txt(s,x+Inches(0.2),yy,colw-Inches(0.4),Inches(0.6),[("定位："+pos,11.5,DARK,True)]);yy=yy+Inches(0.72)
        txt(s,x+Inches(0.2),yy,colw-Inches(0.4),Inches(0.3),[("✔ 优点",12,GREEN,True)]);yy=yy+Inches(0.34)
        for p in pro:
            txt(s,x+Inches(0.3),yy,colw-Inches(0.5),Inches(0.5),[("• "+p,10.5,DARK,False)]);yy=yy+Inches(0.5)
        txt(s,x+Inches(0.2),yy,colw-Inches(0.4),Inches(0.3),[("✘ 缺点",12,ORANGE,True)]);yy=yy+Inches(0.34)
        for c in con:
            txt(s,x+Inches(0.3),yy,colw-Inches(0.5),Inches(0.5),[("• "+c,10.5,GREY,False)]);yy=yy+Inches(0.5)
    fill_card(lx,aws_pos,aws_pro,aws_con)
    fill_card(rx,gcp_pos,gcp_pro,gcp_con)
    # 选型建议条
    by=Inches(6.75)
    box(s,Inches(0.45),by,Inches(12.45),Inches(0.55),fill=DARK,round=True)
    txt(s,Inches(0.7),by,Inches(12.0),Inches(0.55),[("选型建议： "+advice,12.5,WHITE,True)],PP_ALIGN.LEFT,MSO_ANCHOR.MIDDLE)

C=[
(1,"数据仓库","Amazon Redshift","Google BigQuery",
 "PB级托管数仓，Provisioned+Serverless两形态",
 "全托管Serverless、存算分离、AI-ready数据平台",
 ["Serverless无需管集群，按需/预留灵活","与S3/Glue/SageMaker深度集成","RMS存储按GB/月独立计费，成本可控"],
 ["Provisioned仍需选节点/管容量","Python UDF 2026-06-30后停止支持"],
 ["真Serverless、存算彻底分离、资源动态互不干扰","计费双模型可切(按TiB扫描/按slot)","原生AI/ML、搜索、开放表格式内建"],
 ["on-demand按扫描字节，未裁剪大表易失控","深度绑定GCP生态"],
 "重AWS生态/需精细节点控制→Redshift；追求零运维+内建AI+弹性计费→BigQuery"),
(2,"交互式查询","Amazon Athena","BigQuery外部表 / BigLake",
 "Serverless SQL直查S3，另含Serverless Spark",
 "BigQuery外部表/BigLake查GCS及开放格式并统一治理",
 ["完全Serverless、按查询付费、无基础设施","自动并行扩展，大数据集秒级返回","内建Serverless Spark notebook，SQL+Spark双模"],
 ["按扫描量计费，未用列存/分区易贵","高并发/复杂长任务性能受限(待确认配额)"],
 ["一份数据多引擎共享，支持Iceberg/Delta/Hudi","与BigQuery治理(列级安全/脱敏/血缘)统一","可与原生表混查，体验一致"],
 ["外部表相比原生表有性能/功能限制","概念较多(外部表/BigLake/Iceberg)门槛略高"],
 "纯S3即席查询/轻量按需→Athena；已用BQ且要跨引擎共享开放格式→BigLake/外部表"),
(3,"托管Spark/Hadoop","Amazon EMR","Managed Service for Apache Spark",
 "托管Spark/Hive等，EMR on EC2/EKS + EMR Serverless",
 "原Dataproc，托管Spark，1分钟起停、秒级计费",
 ["EMR Serverless自动配资源、跑完即释放","EC2/EKS/Serverless多部署形态灵活","可预初始化资源获秒级交互响应"],
 ["EMR on EC2仍需选实例/调优集群","版本/框架组合与调参复杂度高"],
 ["按vCPU $0.01/小时、秒级计费、1分钟最小计费","集群快速创建、用完关机省钱","支持抢占式实例进一步降本"],
 ["传统集群形态仍需管理生命周期","生态偏Spark，Hadoop组件广度不如EMR"],
 "需多框架/多部署+深AWS集成→EMR；主跑Spark要极简计费快速起停→Managed Spark"),
(4,"流式处理引擎","Kinesis + Managed Flink","Dataflow (Apache Beam)",
 "Kinesis摄取 + Managed Flink做流处理(多语言)",
 "统一批流处理，同一Beam模型，默认exactly-once",
 ["Managed Flink托管底层(弹性/AZ容错/快照)","基于开源Flink，可用全部算子/库","与Kinesis/MSK/S3原生打通"],
 ["需自行组合Streams+Flink两个服务","Flink应用调优/状态管理有学习曲线"],
 ["批流统一同一模型，一套代码两用","全托管、自动分配/回收worker，用完即删","默认exactly-once，可选at-least-once降本"],
 ["绑定Beam模型，跨云迁移成本高","复杂pipeline资源/成本调优需经验"],
 "已用Flink/AWS流栈→Kinesis+Managed Flink；要批流一体极简运维→Dataflow"),
(5,"消息/事件摄取","Kinesis / Amazon MSK","Pub/Sub",
 "Kinesis(原生流摄取) / MSK(托管Kafka)",
 "异步可扩展消息服务，解耦生产消费，延迟约100ms",
 ["MSK提供托管标准Kafka(生态兼容可迁移)","Kinesis与Firehose/Flink/Lambda无缝集成","覆盖Kafka兼容与AWS原生两类需求"],
 ["Kinesis分片需容量规划与扩缩管理","MSK需一定Kafka运维知识"],
 ["全托管、自动扩展、无分片/容量规划负担","异步pub/sub提升系统灵活与健壮","天然对接Dataflow/BigQuery做流分析"],
 ["非Kafka协议，已有Kafka应用迁移需改造","端到端约100ms，超低延迟场景需评估"],
 "需Kafka兼容→MSK；AWS原生流→Kinesis；要零运维全托管事件总线→Pub/Sub"),
(6,"ETL/数据集成","AWS Glue","Cloud Data Fusion",
 "Serverless数据集成(爬取/编目/Spark ETL)",
 "基于开源CDAP的可视化图形化ETL/ELT，丰富插件",
 ["Serverless、无需管理基础设施","内建Data Catalog统一元数据(供Athena/Redshift共享)","与AWS分析栈深度打通"],
 ["可视化编排弱于图形化工具(偏代码)","Spark作业冷启动与调优需经验"],
 ["图形化拖拽建管道、大量预置插件、低代码","基于开源CDAP，可移植性好","支持批/流与多源连接(SAP/Salesforce等)"],
 ["底层跑在Dataproc，实例常驻产生持续费用","相对重，启动/实例管理有开销"],
 "代码化Serverless ETL+AWS编目→Glue；低代码可视化+多源连接→Data Fusion"),
(7,"工作流编排","Amazon MWAA","Managed Airflow (Composer Gen3)",
 "托管Apache Airflow，Python建管道、自动伸缩",
 "全托管Airflow(Gen3)，可跨云及本地编排",
 ["全托管Airflow、自动伸缩、无基础设施负担","建环境时选Airflow版本、自动装配","集成AWS安全服务快速安全访问数据"],
 ["环境按容量常驻计费，空闲也有成本","Airflow版本升级/兼容需规划"],
 ["全托管、原生Airflow Web UI与CLI","可跨云/本地编排pipeline","Gen3持续迭代新特性"],
 ["底层环境常驻、成本随规模上升","版本代际(Gen1/2/3)迁移需注意兼容"],
 "功能相当，按所在云选：AWS→MWAA，GCP→Composer/Managed Airflow"),
(8,"BI/可视化","Amazon Quick Sight","Looker / Looker Studio",
 "交互式BI(现属Amazon Quick AI套件)，可嵌入分析",
 "Looker(LookML语义层企业BI) / Looker Studio(轻量免费)",
 ["Serverless、按会话/容量计费，可嵌入应用","归入Amazon Quick后与AI代理/自然语言打通","与AWS数据源原生连接"],
 ["品牌重组仍在演进，功能边界待稳定","建模/语义层能力弱于Looker"],
 ["LookML语义层统一指标定义、单一真相源","强嵌入分析与API、可作数据应用平台","Looker Studio免费轻量上手快"],
 ["Looker企业版授权成本高、LookML有学习曲线","Looker与Looker Studio两套产品易混淆"],
 "要治理化语义层/企业建模→Looker；轻量免费→Looker Studio；AWS生态+AI化BI→Quick Sight"),
(9,"搜索/日志分析","Amazon OpenSearch","GCP(无原生对标)",
 "托管OpenSearch集群，日志分析/实时监控/点击流",
 "无与OpenSearch一一对应的第一方托管搜索集群",
 ["全托管domain，自动替换故障节点、一键扩缩","兼容OpenSearch与legacy Elasticsearch OSS(≤7.10)","生态含Dashboards、日志/可观测完整"],
 ["domain需选实例/容量并管理","集群规模/分片规划仍需经验"],
 ["Cloud Logging+Log Analytics(底层BigQuery)做检索","全文/向量搜索用BQ search或Vertex AI Search","可在GCP部署Elastic Cloud替代"],
 ["无第一方Elasticsearch兼容托管搜索集群","需自行组合多个服务或用第三方"],
 "需ES/OpenSearch兼容托管集群→AWS明显占优；GCP侧用Log Analytics/BQ+Vertex AI Search或Elastic Cloud"),
(10,"湖仓/表格式","Amazon S3 Tables","BQ Apache Iceberg managed tables",
 "新bucket类型(table bucket)，原生存Iceberg+自动优化",
 "GCS上以Iceberg开放格式提供与原生BQ表一致的托管体验",
 ["原生Iceberg存储+内建自动维护(compaction)","多引擎可查(Athena/Redshift/Spark)","支持Iceberg V3、跨区复制、智能分层降本"],
 ["新bucket类型，区域/配额限制需核对","生态与工具链仍在成熟"],
 ["全托管但数据存客户自有GCS，兼顾开放+治理","支持DML/高吞吐流写/schema演进/时间旅行/列级安全","自动存储优化(自适应文件/聚簇/GC)"],
 ["强绑BigQuery作为管理层","跨引擎写入路径受限(以BQ为主写)"],
 "以S3为湖底座+开放Iceberg多引擎→S3 Tables；已用BQ要托管+数据自有+开放格式→BQ Iceberg managed tables"),
]
for c in C: cat_slide(*c)

# ---------- 总结页 ----------
s=prs.slides.add_slide(blank)
header(s,"总结：怎么选","一句话决策指南")
pts=[
 ("追求零运维、开箱即分析、内建AI","→ GCP(BigQuery为核心)",GREEN),
 ("要灵活组合、生态广度、精细选型","→ AWS(purpose-built服务群)",BLUE),
 ("需要搜索/日志分析(ES兼容)","→ AWS OpenSearch(GCP无原生对标)",ORANGE),
 ("需要托管Kafka","→ AWS MSK(GCP另有Managed Kafka)",ORANGE),
 ("要批流一体、统一编程模型","→ GCP Dataflow(Beam)",GREEN),
 ("要企业级BI语义层","→ GCP Looker(LookML)",GREEN),
 ("要开放中立的托管Iceberg湖","→ AWS S3 Tables(存储归存储,引擎随便挑)",BLUE),
 ("要存算一体的托管Iceberg","→ GCP BQ Iceberg managed tables",GREEN),
]
y=Inches(1.45)
for cond,rec,col in pts:
    box(s,Inches(0.6),y,Inches(6.3),Inches(0.6),fill=LIGHT,line=RGBColor(0xDD,0xE5,0xF0),round=True)
    txt(s,Inches(0.8),y,Inches(6.0),Inches(0.6),[(cond,13,DARK,True)],PP_ALIGN.LEFT,MSO_ANCHOR.MIDDLE)
    box(s,Inches(7.0),y,Inches(5.7),Inches(0.6),fill=col,round=True)
    txt(s,Inches(7.2),y,Inches(5.4),Inches(0.6),[(rec,12.5,WHITE,True)],PP_ALIGN.LEFT,MSO_ANCHOR.MIDDLE)
    y=y+Inches(0.68)

prs.save("/home/ubuntu/.openclaw/workspace/aws_gcp_bigdata_compare.pptx")
print("saved", len(prs.slides.__iter__.__self__._sldIdLst), "slides")
