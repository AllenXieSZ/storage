#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
AWS RDS vs GCP Cloud SQL 对比 PPT
数据来源：AWS RDS User Guide / GCP Cloud SQL 官方文档 + SLA 页，均基于最新官方文档核实。
风格：伟伟认可的 NetApp 风模板（16:9），AWS 橙 / GCP 蓝双列对照。
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

AWS_ORANGE = RGBColor(0xFF, 0x99, 0x00)
GCP_BLUE = RGBColor(0x42, 0x85, 0xF4)
MAIN_BLUE = RGBColor(0x00, 0x67, 0xC5)
DARK = RGBColor(0x1A, 0x2B, 0x4A)
GRAY = RGBColor(0x5A, 0x6B, 0x82)
LIGHT_BG = RGBColor(0xF2, 0xF6, 0xFB)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
ORANGE = RGBColor(0xF5, 0x82, 0x20)
GREEN = RGBColor(0x2E, 0x9E, 0x5B)
RED = RGBColor(0xD0, 0x3A, 0x2E)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
blank = prs.slide_layouts[6]


def add_slide():
    return prs.slides.add_slide(blank)


def bg(slide, color=WHITE):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    s.fill.solid(); s.fill.fore_color.rgb = color
    s.line.fill.background(); s.shadow.inherit = False
    slide.shapes._spTree.remove(s._element)
    slide.shapes._spTree.insert(2, s._element)
    return s


def box(slide, l, t, w, h, fill=None, line=None, line_w=1.0, shape=MSO_SHAPE.RECTANGLE):
    s = slide.shapes.add_shape(shape, l, t, w, h)
    if fill is None:
        s.fill.background()
    else:
        s.fill.solid(); s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line; s.line.width = Pt(line_w)
    s.shadow.inherit = False
    return s


def text(slide, l, t, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, wrap=True):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = Pt(2); tf.margin_right = Pt(2)
    tf.margin_top = Pt(1); tf.margin_bottom = Pt(1)
    if isinstance(runs, str):
        runs = [(runs, 18, DARK, False)]
    first = True
    for item in runs:
        txt, sz, col, bold = item[0], item[1], item[2], item[3]
        sp = item[4] if len(item) > 4 else None
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = align
        if sp is not None:
            p.space_before = Pt(sp)
        r = p.add_run(); r.text = txt
        r.font.size = Pt(sz); r.font.bold = bold
        r.font.color.rgb = col
        r.font.name = "Microsoft YaHei"
    return tb


def title_bar(slide, title, sub=None):
    box(slide, 0, 0, SW, Inches(1.05), fill=MAIN_BLUE)
    box(slide, 0, Inches(1.05), SW, Pt(4), fill=ORANGE)
    text(slide, Inches(0.55), Inches(0.12), Inches(12), Inches(0.8),
         [(title, 26, WHITE, True)], anchor=MSO_ANCHOR.MIDDLE)
    if sub:
        text(slide, Inches(0.6), Inches(1.18), Inches(12.2), Inches(0.4),
             [(sub, 13, GRAY, False)])


def two_col_table(slide, rows, top=Inches(1.75),
                  col_labels=("AWS RDS", "GCP Cloud SQL"),
                  x0=Inches(0.5), label_w=Inches(3.0), col_w=Inches(4.55),
                  row_h=Inches(0.62), fs=12):
    x_aws = x0 + label_w
    x_gcp = x_aws + col_w
    box(slide, x0, top, label_w, Inches(0.5), fill=DARK)
    box(slide, x_aws, top, col_w, Inches(0.5), fill=AWS_ORANGE)
    box(slide, x_gcp, top, col_w, Inches(0.5), fill=GCP_BLUE)
    text(slide, x0, top, label_w, Inches(0.5), [("对比维度", 13, WHITE, True)],
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    text(slide, x_aws, top, col_w, Inches(0.5), [(col_labels[0], 13, WHITE, True)],
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    text(slide, x_gcp, top, col_w, Inches(0.5), [(col_labels[1], 13, WHITE, True)],
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    y = top + Inches(0.5)
    for i, (lab, a, g) in enumerate(rows):
        rb = LIGHT_BG if i % 2 == 0 else WHITE
        box(slide, x0, y, label_w, row_h, fill=RGBColor(0xE8, 0xEE, 0xF6), line=WHITE, line_w=0.5)
        box(slide, x_aws, y, col_w, row_h, fill=rb, line=WHITE, line_w=0.5)
        box(slide, x_gcp, y, col_w, row_h, fill=rb, line=WHITE, line_w=0.5)
        text(slide, x0 + Inches(0.08), y, label_w - Inches(0.12), row_h,
             [(lab, 11.5, MAIN_BLUE, True)], anchor=MSO_ANCHOR.MIDDLE)
        text(slide, x_aws + Inches(0.1), y, col_w - Inches(0.18), row_h,
             [(a, fs, DARK, False)], anchor=MSO_ANCHOR.MIDDLE)
        text(slide, x_gcp + Inches(0.1), y, col_w - Inches(0.18), row_h,
             [(g, fs, DARK, False)], anchor=MSO_ANCHOR.MIDDLE)
        y += row_h


def source_note(slide, txt):
    text(slide, Inches(0.55), Inches(7.02), Inches(12.2), Inches(0.4),
         [(txt, 10, GRAY, False)])


# ============ Slide 1: 封面 ============
s = add_slide(); bg(s, DARK)
box(s, 0, Inches(2.9), SW, Pt(3), fill=ORANGE)
box(s, Inches(0.0), Inches(2.95), SW, Inches(0.02), fill=MAIN_BLUE)
text(s, Inches(0.9), Inches(1.6), Inches(11.5), Inches(1.2),
     [("AWS RDS vs GCP Cloud SQL", 42, WHITE, True)])
text(s, Inches(0.9), Inches(3.15), Inches(11.5), Inches(1.0),
     [("托管关系型数据库全面对比：引擎 · 部署模式 · 存储 · HA/SLA · 特色能力", 19, RGBColor(0x9F,0xC4,0xEA), False)])
text(s, Inches(0.9), Inches(6.5), Inches(11.5), Inches(0.5),
     [("资料来源：AWS RDS User Guide · GCP Cloud SQL 官方文档 · 各自 SLA 页  |  2026-08", 13, GRAY, False)])
box(s, Inches(0.9), Inches(0.7), Inches(1.5), Inches(0.55), fill=AWS_ORANGE)
text(s, Inches(0.9), Inches(0.7), Inches(1.5), Inches(0.55),
     [("AWS RDS", 15, WHITE, True)], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
box(s, Inches(2.55), Inches(0.7), Inches(1.9), Inches(0.55), fill=GCP_BLUE)
text(s, Inches(2.55), Inches(0.7), Inches(1.9), Inches(0.55),
     [("GCP Cloud SQL", 14, WHITE, True)], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


# ============ Slide 2: 对标关系 ============
s = add_slide(); bg(s, WHITE)
title_bar(s, "对标关系：先厘清层次", "本 PPT 聚焦 RDS vs Cloud SQL（托管传统引擎），Aurora/AlloyDB 仅附注")
rows = [
    ("托管传统引擎", "RDS  ← 本对比主角", "Cloud SQL  ← 本对比主角"),
    ("云原生高性能", "Aurora（MySQL/PG 兼容）", "AlloyDB（PostgreSQL 兼容）"),
    ("全球分布式", "Aurora Global / Aurora DSQL", "Spanner"),
]
two_col_table(s, rows, row_h=Inches(0.85), fs=13)
text(s, Inches(0.55), Inches(4.7), Inches(12.2), Inches(1.6),
     [("说明：", 14, MAIN_BLUE, True),
      ("• RDS 与 Cloud SQL 都是「全托管传统数据库引擎」，是最直接的同类对标。", 12.5, DARK, False, 8),
      ("• Aurora / AlloyDB 是各自的云原生高性能引擎，架构是另一档，本次不展开。", 12.5, DARK, False, 4),
      ("• 全球分布式场景（Spanner / Aurora Global）不属于同类，仅列出对应关系。", 12.5, DARK, False, 4)])
source_note(s, "来源：docs.aws.amazon.com/AmazonRDS · cloud.google.com/sql")


# ============ Slide 3: 支持引擎 ============
s = add_slide(); bg(s, WHITE)
title_bar(s, "1. 支持的数据库引擎", "RDS 引擎覆盖更广，独有 Oracle / Db2 / MariaDB")
rows = [
    ("引擎数量", "6 种", "3 种"),
    ("MySQL / PostgreSQL", "✅ / ✅", "✅ / ✅"),
    ("SQL Server", "✅", "✅"),
    ("MariaDB", "✅", "❌"),
    ("Oracle", "✅（全托管）", "❌（走 Oracle@Google Cloud）"),
    ("IBM Db2", "✅（全托管）", "❌（无对标托管服务）"),
]
two_col_table(s, rows, row_h=Inches(0.65))
source_note(s, "来源：docs.aws.amazon.com/AmazonRDS/latest/UserGuide/Welcome.html · developers.google.com/cloud-sql")


# ============ Slide 4: GCP 上跑 Oracle / Db2 ============
s = add_slide(); bg(s, WHITE)
title_bar(s, "2. GCP 上如何跑 Oracle / Db2？", "Cloud SQL 不支持，需走其它路径")
# Oracle 卡片
box(s, Inches(0.55), Inches(1.7), Inches(6.05), Inches(4.5), fill=LIGHT_BG,
    line=GCP_BLUE, line_w=1.5, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
text(s, Inches(0.8), Inches(1.9), Inches(5.6), Inches(0.5),
     [("Oracle → 有官方方案", 17, GCP_BLUE, True)])
text(s, Inches(0.8), Inches(2.5), Inches(5.6), Inches(3.5),
     [("Oracle Database@Google Cloud", 14, DARK, True),
      ("Google × Oracle 合作：把 OCI Exadata 硬件", 12, DARK, False, 8),
      ("放进 GCP 数据中心，低延迟直连 + IAM 集成。", 12, DARK, False, 2),
      ("", 6, DARK, False),
      ("支持的 OCI 服务：", 12.5, MAIN_BLUE, True),
      ("• Exadata Database Service", 12, DARK, False, 4),
      ("• Autonomous AI Database Service", 12, DARK, False, 2),
      ("• Base Database Service", 12, DARK, False, 2),
      ("• Exascale Infrastructure / GoldenGate", 12, DARK, False, 2),
      ("", 6, DARK, False),
      ("本质：不是 Cloud SQL 托管，是在 GCP 里", 11.5, GRAY, False),
      ("用原生 OCI Oracle 服务。", 11.5, GRAY, False, 2)])
# Db2 卡片
box(s, Inches(6.85), Inches(1.7), Inches(6.05), Inches(4.5), fill=RGBColor(0xFB,0xEC,0xEC),
    line=RED, line_w=1.5, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
text(s, Inches(7.1), Inches(1.9), Inches(5.6), Inches(0.5),
     [("Db2 → 无对标托管", 17, RED, True)])
text(s, Inches(7.1), Inches(2.5), Inches(5.6), Inches(3.5),
     [("GCP 没有等价的全托管 Db2 服务", 14, DARK, True),
      ("（Cloud SQL 不支持，也无 Oracle 那样的合作方案）", 11.5, DARK, False, 4),
      ("", 6, DARK, False),
      ("客户只能：", 12.5, MAIN_BLUE, True),
      ("① Compute Engine 上自建 Db2（非托管）", 12, DARK, False, 4),
      ("② IBM Cloud 的 Db2 托管（跨云）", 12, DARK, False, 2),
      ("③ 迁移到 PostgreSQL / 其它引擎", 12, DARK, False, 2),
      ("", 6, DARK, False),
      ("→ 这是 GCP 相对 AWS 的真实短板：", 12, RED, True),
      ("   AWS RDS 有全托管 Db2。", 12, RED, True, 2)])
source_note(s, "来源：cloud.google.com/oracle/database/docs/overview（Oracle@Google Cloud）")


# ============ Slide 5: 版本模型（edition）============
s = add_slide(); bg(s, WHITE)
title_bar(s, "3. 版本模型：Cloud SQL 分 Edition", "Cloud SQL 两档打包能力；RDS 靠实例类+选项自由组合")
rows = [
    ("模型", "无 edition 概念，靠实例类+选项组合", "Enterprise / Enterprise Plus 两档"),
    ("可用性 SLA", "统一（见 SLA 页）", "Ent 99.95% / Ent Plus 99.99%"),
    ("Data cache（本地SSD读缓存）", "对应 Optimized Reads（本地NVMe）", "仅 Enterprise Plus，读性能最高 4x"),
    ("优化写", "Optimized Writes（NVMe 16KiB 原子写）", "仅 Ent Plus，写吞吐 3x / 延迟-98%"),
    ("PITR 日志保留", "0–35 天", "Ent 最多 7 天 / Ent Plus 最多 35 天"),
    ("计划维护停机", "维护窗口内（可 Blue/Green 近零）", "Ent 数分钟 / Ent Plus 亚秒级"),
]
two_col_table(s, rows, row_h=Inches(0.65), fs=11)
source_note(s, "来源：docs.cloud.google.com/sql/docs/mysql/editions-intro（官方 edition 对比表）")


# ============ Slide 6: RDS 三种部署模式 ============
s = add_slide(); bg(s, WHITE)
title_bar(s, "4. RDS 三种部署模式（含本地 NVMe binlog）", "Multi-AZ DB cluster 用本地 NVMe 存事务日志/binlog 降写延迟")
modes = [
    ("① Single-AZ instance", "单实例，单 AZ", "开发 / 测试，无高可用", GRAY),
    ("② Multi-AZ DB instance", "1 主 + 1 备用(standby, 不可读)，2 AZ，同步复制", "高可用，故障自动切换，备用不分担读", MAIN_BLUE),
    ("③ Multi-AZ DB cluster", "1 写 + 2 可读 reader，跨 3 AZ，半同步复制", "HA + 读扩展 + 更低写延迟", GREEN),
]
y = Inches(1.7)
for t, arch, use, col in modes:
    box(s, Inches(0.55), y, Inches(0.16), Inches(1.0), fill=col)
    text(s, Inches(0.85), y + Inches(0.05), Inches(11.9), Inches(0.4), [(t, 15, col, True)])
    text(s, Inches(0.85), y + Inches(0.42), Inches(11.9), Inches(0.3), [("架构：" + arch, 12, DARK, False)])
    text(s, Inches(0.85), y + Inches(0.72), Inches(11.9), Inches(0.3), [("用途：" + use, 11.5, GRAY, False)])
    y += Inches(1.08)
# 高亮框
box(s, Inches(0.55), Inches(5.05), Inches(12.25), Inches(1.5), fill=RGBColor(0xEC,0xF7,0xF0),
    line=GREEN, line_w=1.2, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
text(s, Inches(0.75), Inches(5.15), Inches(11.9), Inches(1.4),
     [("关键：③ Multi-AZ DB cluster「binlog/事务日志写本地 NVMe」", 13.5, GREEN, True),
      ("• 强制使用带 d 后缀的本地 NVMe 实例类（db.m6gd/m6id/r6gd/r5d/x2iedn…）", 11.5, DARK, False, 5),
      ("• RDS 官方特性页原文：\"Uses local storage for transactional logs to reduce jitter\"", 11.5, DARK, False, 3),
      ("• binlog 走本地 NVMe(快) + 半同步复制到 2 个 AZ reader → 写延迟低于 Multi-AZ DB instance", 11.5, DARK, False, 3)])
source_note(s, "来源：multi-az-db-clusters-concepts.html · aws.amazon.com/rds/features/multi-az · USER_Binlog.MultiAZ.html")


# ============ Slide 7: 高可用 & SLA ============
s = add_slide(); bg(s, WHITE)
title_bar(s, "5. 高可用 & SLA", "HA 基础档打平；Cloud SQL Enterprise Plus 高一档；RDS 单实例有兜底")
rows = [
    ("HA / 多可用区 SLA", "99.95%（Multi-AZ 实例 / 集群）", "99.95%（Ent）/ 99.99%（Ent Plus）"),
    ("单实例 SLA", "99.5%（Single-DB Instance 有兜底）", "单区/共享核实例不在 SLA 覆盖内"),
    ("HA 机制", "同步复制到备用副本，自动故障转移", "跨 zone 主备，故障自动切换"),
    ("读扩展 HA", "Multi-AZ DB cluster（2 可读 reader）", "Read pool（Ent Plus，1–20 节点）"),
]
two_col_table(s, rows, row_h=Inches(0.8), fs=12)
source_note(s, "来源：aws.amazon.com/rds/sla · cloud.google.com/sql/sla · about-read-pools")


# ============ Slide 8: 存储 & 副本 ============
s = add_slide(); bg(s, WHITE)
title_bar(s, "6. 存储上限 & 只读副本", "RDS 存储可调性更强；副本模型两家不同")
rows = [
    ("最大存储", "64 TiB；Oracle/SQLServer 附加卷 256 TiB", "64 TB（专用核）/ 3 TB（共享核）"),
    ("持久块存储", "gp3 / io2 Block Express / io1", "SSD / HDD"),
    ("最高 IOPS", "io2 Block Express 256,000 IOPS", "随规格自动走，无独立 provisioned IOPS"),
    ("本地 NVMe", "d 系实例类(≤7.6TB) + Optimized Reads", "Enterprise Plus data cache（本地 SSD）"),
    ("只读副本上限", "每主实例最多 15 个（硬上限）", "副本=实例，受项目配额≤1000/项目"),
]
two_col_table(s, rows, row_h=Inches(0.72), fs=11)
source_note(s, "来源：CHAP_Storage.html · Concepts.DBInstanceClass.Types.html · USER_ReadRepl.html · cloud.google.com/sql/docs/quotas")


# ============ Slide 9: 备份 & 特色能力 ============
s = add_slide(); bg(s, WHITE)
title_bar(s, "7. 备份 / PITR & 特色能力", "各自亮点")
# 左 RDS
box(s, Inches(0.55), Inches(1.7), Inches(6.05), Inches(4.9), fill=LIGHT_BG,
    line=AWS_ORANGE, line_w=1.5, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
text(s, Inches(0.8), Inches(1.85), Inches(5.6), Inches(0.4), [("AWS RDS 特色", 16, AWS_ORANGE, True)])
text(s, Inches(0.8), Inches(2.4), Inches(5.6), Inches(4.0),
     [("• 自动备份 0–35 天 + PITR 秒级恢复点", 12.5, DARK, False),
      ("• Blue/Green 部署（近零停机升级/切换）", 12.5, DARK, False, 6),
      ("• RDS Proxy（连接池、故障转移加速）", 12.5, DARK, False, 6),
      ("• Performance Insights（性能分析）", 12.5, DARK, False, 6),
      ("• RDS Custom（可访问 OS/DB，Oracle/SQLServer）", 12.5, DARK, False, 6),
      ("• 引擎最全（6 种，含 Oracle/Db2）", 12.5, DARK, False, 6),
      ("• Optimized Reads / Optimized Writes（本地 NVMe）", 12.5, DARK, False, 6)])
# 右 Cloud SQL
box(s, Inches(6.85), Inches(1.7), Inches(6.05), Inches(4.9), fill=LIGHT_BG,
    line=GCP_BLUE, line_w=1.5, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
text(s, Inches(7.1), Inches(1.85), Inches(5.6), Inches(0.4), [("GCP Cloud SQL 特色", 16, GCP_BLUE, True)])
text(s, Inches(7.1), Inches(2.4), Inches(5.6), Inches(4.0),
     [("• PITR：Ent 7 天 / Ent Plus 35 天", 12.5, DARK, False),
      ("• near-zero downtime 计划维护（Ent Plus 亚秒）", 12.5, DARK, False, 6),
      ("• Managed Connection Pooling（Ent Plus）", 12.5, DARK, False, 6),
      ("• Query Insights + AI 辅助排障（Ent Plus）", 12.5, DARK, False, 6),
      ("• IAM 数据库认证深度集成 + Auth Proxy", 12.5, DARK, False, 6),
      ("• Read pools（Ent Plus，1–20 节点统一读端点）", 12.5, DARK, False, 6),
      ("• Data cache 读 4x / 优化写 3x（Ent Plus）", 12.5, DARK, False, 6)])
source_note(s, "来源：各自官方 feature 文档 · editions-intro")


# ============ Slide 10: 云原生篇 分隔页 ============
s = add_slide(); bg(s, DARK)
box(s, 0, Inches(2.9), SW, Pt(3), fill=ORANGE)
text(s, Inches(0.9), Inches(2.0), Inches(11.5), Inches(1.0),
     [("云原生篇：Aurora vs AlloyDB", 40, WHITE, True)])
text(s, Inches(0.9), Inches(3.3), Inches(11.5), Inches(1.0),
     [("存算分离的云原生数据库对标（RDS/Cloud SQL 之上的一档）", 19, RGBColor(0x9F,0xC4,0xEA), False)])
box(s, Inches(0.9), Inches(4.6), Inches(1.5), Inches(0.5), fill=AWS_ORANGE)
text(s, Inches(0.9), Inches(4.6), Inches(1.5), Inches(0.5),
     [("Aurora", 14, WHITE, True)], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
box(s, Inches(2.55), Inches(4.6), Inches(1.5), Inches(0.5), fill=GCP_BLUE)
text(s, Inches(2.55), Inches(4.6), Inches(1.5), Inches(0.5),
     [("AlloyDB", 14, WHITE, True)], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


# ============ Slide 11: Aurora vs AlloyDB 架构 ============
s = add_slide(); bg(s, WHITE)
title_bar(s, "8. 云原生架构：存算分离", "两者都是「基于 LOG 的共享存储 + 读写分离」（墨天轮解析印证）")
rows = [
    ("兼容引擎", "MySQL + PostgreSQL（双引擎）", "仅 PostgreSQL"),
    ("存储模型", "cluster volume：数据 6 副本跨 3 AZ", "disaggregated 存算解耦，跨多 AZ"),
    ("存储分层", "分布式共享存储卷（SSD）", "log storage + LPS(日志处理) + 块存储 三层"),
    ("binlog/日志", "binlog 存在 cluster volume 内", "WAL 转日志写入 log storage 层"),
    ("只读副本", "最多 15 个 Aurora Replica，共享存储", "主节点 + read pool，共享 LOG 存储"),
    ("临时文件", "本地存储放排序/建索引临时文件", "计算节点本地资源"),
]
two_col_table(s, rows, col_labels=("AWS Aurora", "GCP AlloyDB"), row_h=Inches(0.62), fs=11)
source_note(s, "来源：Aurora.Overview.StorageReliability.html(官方) · alloydb/docs/overview(官方) · 三层存储与LOG架构：墨天轮 modb.pro/db/445739")


# ============ Slide 12: HTAP / 列存 / AI ============
s = add_slide(); bg(s, WHITE)
title_bar(s, "9. HTAP / 列式引擎 / AI（AlloyDB 差异化）", "AlloyDB 内置列存内存 + WAL Apply 优化，主打 HTAP")
rows = [
    ("HTAP / 列式", "本体偏 OLTP，分析走 Redshift/zero-ETL", "内置 columnar engine（列存内存），原生 HTAP"),
    ("列存装载方式", "—", "手工装载：SELECT google_columnar_engine_add('表')"),
    ("只读节点优化", "共享存储读", "WAL Apply 优化：只读节点仅重演主实例\nCACHE 数据，其余直接读数据文件"),
    ("AI 能力", "pgvector 扩展 + Bedrock/SageMaker 集成", "AlloyDB AI：向量搜索+ML 调用内置进引擎"),
    ("实现方式", "AWS 自研引擎", "通过 PG Extension 实现，随 PG 社区版演进"),
]
two_col_table(s, rows, col_labels=("AWS Aurora", "GCP AlloyDB"), row_h=Inches(0.85), fs=11)
source_note(s, "来源：alloydb/docs/overview(官方 columnar/AlloyDB AI) · 列存装载/WAL Apply优化/Extension实现：墨天轮 modb.pro/db/445739")


# ============ Slide 13: 命名 + 性能宣称 + 结论 ============
s = add_slide(); bg(s, WHITE)
title_bar(s, "10. AlloyDB 命名 & 性能宣称（分层标注）", "厂商 benchmark 互不承认，命名无官方逐字解释")
# 命名卡片
box(s, Inches(0.55), Inches(1.65), Inches(12.25), Inches(1.5), fill=LIGHT_BG,
    line=GCP_BLUE, line_w=1.2, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
text(s, Inches(0.75), Inches(1.75), Inches(11.9), Inches(1.4),
     [("「Alloy」= 合金：为什么这么取名？", 15, GCP_BLUE, True),
      ("Alloy(合金)=把多种金属熔合成性能更强的新材料。AlloyDB 寓意把「开源 PostgreSQL」熔合「Google 自研 scale-out 计算/存储 + AI/ML」= 更强的新产品。", 12, DARK, False, 5),
      ("⚠️ Google 无「官方逐字」命名解释文档；此为基于官方定位(blog: \"combines the best of Google's... with full PostgreSQL compatibility\")的合理推断。", 10.5, GRAY, False, 4)])
# 性能宣称
box(s, Inches(0.55), Inches(3.3), Inches(12.25), Inches(1.35), fill=RGBColor(0xFF,0xF6,0xE8),
    line=ORANGE, line_w=1.2, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
text(s, Inches(0.75), Inches(3.4), Inches(11.9), Inches(1.25),
     [("性能宣称（⚠️ 均为厂商自家 benchmark，互不承认，选型须自测）", 14, ORANGE, True),
      ("• AlloyDB(GCP自述)：比原生 PG 事务 4x、分析 100x；比 Amazon 同类事务 2x   • Aurora(AWS自述)：比标准 MySQL 5x、PG 3x", 11.5, DARK, False, 5),
      ("• 第三方(墨天轮)：AlloyDB 大型统计 SQL 实测提升几十倍，但非严格对照；AlloyDB 未像 HeatWave 那样堆硬件砸性能", 11, GRAY, False, 3)])
# 结论
items = [
    ("引擎", "要 MySQL 云原生只能 Aurora；AlloyDB 只做 PostgreSQL", AWS_ORANGE),
    ("架构", "都是 LOG 共享存储+读写分离；Aurora 6副本/3AZ，AlloyDB 存算解耦三层", GREEN),
    ("差异化", "AlloyDB 内置列存 HTAP + AlloyDB AI（向量/ML内置引擎）更激进", GCP_BLUE),
]
y = Inches(4.85)
for t, body, col in items:
    box(s, Inches(0.55), y, Inches(0.14), Inches(0.62), fill=col)
    text(s, Inches(0.8), y, Inches(11.9), Inches(0.35), [(t, 13.5, col, True)])
    text(s, Inches(0.8), y + Inches(0.32), Inches(11.9), Inches(0.35), [(body, 11.5, DARK, False)])
    y += Inches(0.68)
source_note(s, "命名为推断(Google无官方定义)；性能均厂商自述benchmark；第三方=墨天轮 modb.pro/db/445739")


# ============ Slide 14: 全球分布式篇 分隔页 ============
s = add_slide(); bg(s, DARK)
box(s, 0, Inches(2.9), SW, Pt(3), fill=ORANGE)
text(s, Inches(0.9), Inches(2.0), Inches(11.5), Inches(1.0),
     [("全球分布式篇：Aurora Global / DSQL vs Spanner", 32, WHITE, True)])
text(s, Inches(0.9), Inches(3.3), Inches(11.5), Inches(1.0),
     [("跨区域分布式数据库对标（云原生之上的一档）", 19, RGBColor(0x9F,0xC4,0xEA), False)])
box(s, Inches(0.9), Inches(4.6), Inches(2.4), Inches(0.5), fill=AWS_ORANGE)
text(s, Inches(0.9), Inches(4.6), Inches(2.4), Inches(0.5),
     [("Aurora Global / DSQL", 12, WHITE, True)], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
box(s, Inches(3.45), Inches(4.6), Inches(1.8), Inches(0.5), fill=GCP_BLUE)
text(s, Inches(3.45), Inches(4.6), Inches(1.8), Inches(0.5),
     [("Cloud Spanner", 12, WHITE, True)], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


# ============ Slide 15: 三方对标 + 架构/写模式 ============
s = add_slide(); bg(s, WHITE)
title_bar(s, "11. 全球分布式：三方对标", "真正对标 Spanner 的是 Aurora DSQL（active-active），不是 Aurora Global（单主）")
# 三列表
x0 = Inches(0.5); c1 = Inches(4.1); cw = Inches(4.35)
xa = x0 + c1; xb = xa + cw
top = Inches(1.75)
box(s, x0, top, c1, Inches(0.5), fill=DARK)
box(s, xa, top, cw, Inches(0.5), fill=AWS_ORANGE)
box(s, xb, top, cw, Inches(0.5), fill=GCP_BLUE)
text(s, x0, top, c1, Inches(0.5), [("维度", 12, WHITE, True)], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
text(s, xa, top, cw, Inches(0.5), [("AWS Aurora Global / DSQL", 11.5, WHITE, True)], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
text(s, xb, top, cw, Inches(0.5), [("GCP Cloud Spanner", 12, WHITE, True)], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
rows3 = [
    ("定位", "Global=跨区灾备+就近读\nDSQL=全球分布式 serverless SQL", "全球分布式 SQL（多模型）"),
    ("写模式", "Global=单主(1区写)\nDSQL=active-active 多区可写", "多区，Paxos leader 写"),
    ("兼容", "MySQL/PG（Global）\nPG 兼容（DSQL）", "专有 + PostgreSQL 接口"),
    ("复制", "Global=存储级异步 <1s\nDSQL=同步强一致", "同步 Paxos，多数派提交"),
    ("RPO/RTO", "Global: RPO<1s, RTO<1min\nDSQL: RPO=0 自动恢复", "RPO=0，Paxos 多数派容错"),
]
y = top + Inches(0.5)
for i, (lab, a, g) in enumerate(rows3):
    rb = LIGHT_BG if i % 2 == 0 else WHITE
    rh = Inches(0.82)
    box(s, x0, y, c1, rh, fill=RGBColor(0xE8,0xEE,0xF6), line=WHITE, line_w=0.5)
    box(s, xa, y, cw, rh, fill=rb, line=WHITE, line_w=0.5)
    box(s, xb, y, cw, rh, fill=rb, line=WHITE, line_w=0.5)
    text(s, x0+Inches(0.08), y, c1-Inches(0.12), rh, [(lab, 11.5, MAIN_BLUE, True)], anchor=MSO_ANCHOR.MIDDLE)
    text(s, xa+Inches(0.1), y, cw-Inches(0.18), rh, [(a, 10, DARK, False)], anchor=MSO_ANCHOR.MIDDLE)
    text(s, xb+Inches(0.1), y, cw-Inches(0.18), rh, [(g, 10.5, DARK, False)], anchor=MSO_ANCHOR.MIDDLE)
    y += rh
source_note(s, "来源：aws.amazon.com/rds/aurora/global-database · docs.aws.amazon.com/aurora-dsql · spanner/docs/replication（均官方）")


# ============ Slide 16: 一致性 & SLA ============
s = add_slide(); bg(s, WHITE)
title_bar(s, "12. 一致性机制 & SLA", "Spanner 靠 TrueTime，DSQL 靠 OCC；多区 SLA 都是 99.999%")
rows = [
    ("一致性模型", "DSQL：强一致+snapshot isolation(OCC)\nGlobal：主区强一致/从区<1s滞后", "external consistency（最强）+ TrueTime"),
    ("时钟/并发", "DSQL：QP 互不通信，OCC 乐观并发", "TrueTime(GPS+原子钟) + MVCC + Paxos"),
    ("多区 SLA", "DSQL 多区 99.999% / 单区 99.99%", "Multi/Dual-Regional 99.999% / Regional 99.99%"),
    ("多模型", "纯关系型(PG 兼容)", "关系+Graph+全文+向量 一库多模"),
    ("成熟度", "DSQL 2024 re:Invent 新出，生态较新", "2017 GA，成熟"),
]
two_col_table(s, rows, col_labels=("AWS Aurora DSQL/Global", "GCP Cloud Spanner"),
              label_w=Inches(2.7), col_w=Inches(4.7), row_h=Inches(0.85), fs=10.5)
source_note(s, "来源：docs.aws.amazon.com/aurora-dsql · spanner/sla · spanner true-time-external-consistency（官方）· DSQL OCC/QP：Marc Brooker(AWS DSQL首席工程师)博客 brooker.co.za")


# ============ Slide 17: Spanner 命名 + CAP + latency ============
s = add_slide(); bg(s, WHITE)
title_bar(s, "13. Spanner 命名 & 强一致的代价", "全球强一致牺牲了什么？CAP 与延迟 trade-off")
# 命名
box(s, Inches(0.55), Inches(1.62), Inches(12.25), Inches(1.05), fill=LIGHT_BG,
    line=GCP_BLUE, line_w=1.2, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
text(s, Inches(0.75), Inches(1.7), Inches(11.9), Inches(1.0),
     [("「Spanner」命名来头", 14, GCP_BLUE, True),
      ("寓意「span 跨越」——一个数据库横跨全球多个数据中心/区域（global database that spans the world）。", 11.5, DARK, False, 4),
      ("⚠️ Google 无官方命名解释；此为社区共识（Wikipedia/Reddit），非官方逐字定义。", 10, GRAY, False, 3)])
# CAP
box(s, Inches(0.55), Inches(2.82), Inches(12.25), Inches(1.55), fill=RGBColor(0xEC,0xF7,0xF0),
    line=GREEN, line_w=1.2, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
text(s, Inches(0.75), Inches(2.9), Inches(11.9), Inches(1.5),
     [("CAP：Spanner 技术上是 CP 系统（牺牲 A，非 CA）", 14, GREEN, True),
      ("Eric Brewer(CAP定理提出者/Google VP)官方博客原话：\"during some partitions, Spanner chooses C and forfeits A. It is technically a CP system.\"", 11, DARK, False, 4),
      ("但靠 Google 私有全球光纤网，网络分区极罕见(网络仅占<10%故障)→实测 >5个9(99.999%)，用户「当它是 CA」用。", 11, DARK, False, 3)])
# latency
box(s, Inches(0.55), Inches(4.52), Inches(12.25), Inches(1.7), fill=RGBColor(0xFF,0xF6,0xE8),
    line=ORANGE, line_w=1.2, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
text(s, Inches(0.75), Inches(4.6), Inches(11.9), Inches(1.65),
     [("Latency：牺牲的是「写延迟」（读靠 MVCC 基本不受影响）", 14, ORANGE, True),
      ("① TrueTime commit-wait：写事务要等时钟不确定性区间 ε 过去才提交（保证 external consistency）→ 写多等 ε(毫秒级)", 11, DARK, False, 4),
      ("② 跨区写要 Paxos 多数派同步 → 写延迟受跨区网络 RTT/光速制约（跨洲几十~上百 ms）", 11, DARK, False, 3),
      ("③ 读：MVCC 多版本，snapshot read 不阻塞写、可本地副本完成 → 读延迟低。核心 trade-off = 写延迟↑ 换 全球强一致+高可用", 11, DARK, False, 3)])
source_note(s, "来源：Google Cloud 官方博客《Inside Cloud Spanner and the CAP Theorem》(Eric Brewer) · spanner true-time 官方文档 · 命名：Wikipedia/社区推测")


# ============ Slide 18: 总结 ============
s = add_slide(); bg(s, WHITE)
title_bar(s, "总结：关键差异一览", "选型时优先看这几点")
items = [
    ("引擎覆盖", "RDS 完胜（6 vs 3，独有 Oracle/Db2/MariaDB）；GCP Oracle 靠 Oracle@GCP，Db2 无对标", AWS_ORANGE),
    ("HA / SLA", "HA 基础档打平 99.95%；Cloud SQL Ent Plus 99.99% 高一档；RDS 单实例 99.5% 兜底", GCP_BLUE),
    ("存储可调性", "RDS 更强：io2 Block Express 256K IOPS、256TiB 附加卷、独立 provisioned IOPS", AWS_ORANGE),
    ("部署模式", "RDS 三档；Multi-AZ DB cluster 用本地 NVMe 存 binlog/事务日志降写延迟（强制 d 系）", GREEN),
    ("版本模型", "Cloud SQL 用 Enterprise/Ent Plus 分档打包能力；RDS 靠实例类+选项自由组合", GCP_BLUE),
    ("本地 NVMe", "两家都重度用本地 NVMe/SSD 加速，但用途不同：RDS 临时对象/原子写 vs Cloud SQL 读缓存", GREEN),
    ("全球分布式", "对标 Spanner 的是 Aurora DSQL(active-active)；Spanner=CP系统靠TrueTime，写延迟换全球强一致", GCP_BLUE),
]
y = Inches(1.55)
for title_t, body, col in items:
    box(s, Inches(0.55), y, Inches(0.15), Inches(0.7), fill=col)
    text(s, Inches(0.85), y, Inches(11.9), Inches(0.38), [(title_t, 14, col, True)])
    text(s, Inches(0.85), y + Inches(0.34), Inches(11.9), Inches(0.42), [(body, 11.5, DARK, False)])
    y += Inches(0.77)
source_note(s, "所有数据基于 AWS RDS User Guide / GCP Cloud SQL 官方文档核实；具体以各自 SLA/定价/quotas 页实时为准。")

prs.save("RDS_vs_CloudSQL.pptx")
print("saved RDS_vs_CloudSQL.pptx")
