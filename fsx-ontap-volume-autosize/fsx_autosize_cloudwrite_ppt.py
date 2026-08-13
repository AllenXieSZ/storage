#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
FSx ONTAP Autosizing + Cloud Write Mode PPT
基于 AWS FSx ONTAP 官方文档
NetApp 风格（沿用 TOOLS.md PPT 规范）
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

NETAPP_BLUE = RGBColor(0x00, 0x67, 0xC5)
DARK = RGBColor(0x1A, 0x2B, 0x4A)
GRAY = RGBColor(0x5A, 0x6B, 0x82)
LIGHT_BG = RGBColor(0xF2, 0xF6, 0xFB)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
ORANGE = RGBColor(0xF5, 0x82, 0x20)
GREEN = RGBColor(0x2E, 0x9E, 0x5B)
ACCENT = RGBColor(0x0B, 0x5C, 0xAB)
CLOUD = RGBColor(0x6A, 0x5A, 0xCD)
CODEBG = RGBColor(0x1E, 0x2A, 0x38)
CODEFG = RGBColor(0xE6, 0xED, 0xF3)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
blank = prs.slide_layouts[6]


def add_slide(): return prs.slides.add_slide(blank)


def bg(slide, color=WHITE):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    s.fill.solid(); s.fill.fore_color.rgb = color
    s.line.fill.background(); s.shadow.inherit = False
    slide.shapes._spTree.remove(s._element)
    slide.shapes._spTree.insert(2, s._element)
    return s


def box(slide, l, t, w, h, fill=None, line=None, line_w=1.0, shape=MSO_SHAPE.RECTANGLE):
    s = slide.shapes.add_shape(shape, l, t, w, h)
    if fill is None: s.fill.background()
    else: s.fill.solid(); s.fill.fore_color.rgb = fill
    if line is None: s.line.fill.background()
    else: s.line.color.rgb = line; s.line.width = Pt(line_w)
    s.shadow.inherit = False
    return s


def text(slide, l, t, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, wrap=True, font="Microsoft YaHei"):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = wrap; tf.vertical_anchor = anchor
    tf.margin_left = Pt(2); tf.margin_right = Pt(2); tf.margin_top = Pt(1); tf.margin_bottom = Pt(1)
    if isinstance(runs, str): runs = [(runs, 18, DARK, False)]
    first = True
    for item in runs:
        txt, sz, col, bold = item[0], item[1], item[2], item[3]
        sp = item[4] if len(item) > 4 else None
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = align
        if sp is not None: p.space_before = Pt(sp)
        r = p.add_run(); r.text = txt
        r.font.size = Pt(sz); r.font.bold = bold; r.font.color.rgb = col
        r.font.name = font
    return tb


def title_bar(slide, title, sub=None):
    box(slide, 0, 0, SW, Inches(1.05), fill=NETAPP_BLUE)
    box(slide, 0, Inches(1.05), SW, Pt(4), fill=ORANGE)
    text(slide, Inches(0.55), Inches(0.12), Inches(12), Inches(0.8),
         [(title, 28, WHITE, True)], anchor=MSO_ANCHOR.MIDDLE)
    if sub:
        text(slide, Inches(0.6), Inches(1.2), Inches(12), Inches(0.4),
             [(sub, 14, GRAY, False)])


def codebox(slide, l, t, w, h, lines):
    box(slide, l, t, w, h, fill=CODEBG, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    runs = []
    for i, ln in enumerate(lines):
        col = ORANGE if ln.strip().startswith("#") else CODEFG
        runs.append((ln, 11, col, False, 2 if i else 0))
    text(slide, l+Inches(0.2), t+Inches(0.12), w-Inches(0.4), h-Inches(0.24),
         runs, font="Consolas")


# ============ Slide 1: 封面 ============
s = add_slide(); bg(s, DARK)
box(s, 0, Inches(2.9), SW, Pt(3), fill=ORANGE)
box(s, 0, Inches(2.95), SW, Inches(0.02), fill=NETAPP_BLUE)
text(s, Inches(0.9), Inches(1.65), Inches(11.5), Inches(1.2),
     [("FSx ONTAP 卷管理特性", 42, WHITE, True)])
text(s, Inches(0.9), Inches(3.15), Inches(11.5), Inches(1.0),
     [("Volume Autosizing 自动扩缩容  +  Cloud Write Mode 云直写", 22, RGBColor(0x9F,0xC4,0xEA), False)])
text(s, Inches(0.9), Inches(6.5), Inches(11.5), Inches(0.5),
     [("资料来源：AWS FSx for ONTAP 官方文档 docs.aws.amazon.com  |  2026-06", 13, GRAY, False)])
box(s, Inches(0.9), Inches(0.7), Inches(2.2), Inches(0.55), fill=NETAPP_BLUE)
text(s, Inches(0.9), Inches(0.7), Inches(2.2), Inches(0.55),
     [("NetApp ONTAP", 16, WHITE, True)], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


# ============ Slide 2: Autosizing 原理 ============
s = add_slide(); bg(s, WHITE)
title_bar(s, "Autosizing 自动扩缩容 — 原理", "卷根据「已用空间百分比」自动 grow / shrink")

lx = Inches(0.7)
box(s, lx, Inches(1.45), Inches(6.0), Inches(5.4), fill=LIGHT_BG, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
text(s, lx+Inches(0.3), Inches(1.6), Inches(5.4), Inches(0.5),
     [("是什么", 20, NETAPP_BLUE, True)])
text(s, lx+Inches(0.3), Inches(2.25), Inches(5.4), Inches(4.5),
     [("• 卷的容量随「已用空间」自动调整", 14, DARK, False),
      ("• 仅适用于 FlexVol（FSx 默认卷类型）", 14, DARK, False, 8),
      ("• 两种模式：", 14, DARK, True, 10),
      ("    grow —— 只增不减", 13, GREEN, True, 4),
      ("    grow_shrink —— 可增可减", 13, GREEN, True, 2),
      ("• grow 触发：已用 ≥ grow-threshold（如90%）", 14, DARK, False, 10),
      ("   → 卷自动增大（不超过 max-size）", 12, GRAY, False, 2),
      ("• shrink 触发：已用 ≤ shrink-threshold", 14, DARK, False, 8),
      ("   → 卷自动缩小（不低于 min-size）", 12, GRAY, False, 2),
      ("• max-size 上限 300TB；默认为卷大小的120%", 14, DARK, False, 10),
      ("• 价值：避免卷写满中断，又不浪费预留", 14, ORANGE, True, 10)])

# 右侧示意：阈值触发增长
rx = Inches(7.0)
box(s, rx, Inches(1.45), Inches(5.65), Inches(5.4), fill=WHITE, line=NETAPP_BLUE, line_w=1.5, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
text(s, rx+Inches(0.25), Inches(1.55), Inches(5.1), Inches(0.4),
     [("grow 触发示意", 16, NETAPP_BLUE, True)])
# 卷1：未到阈值
def volbar(x, label, usedpct, col, note):
    box(s, x, Inches(2.3), Inches(1.4), Inches(3.6), fill=RGBColor(0xE8,0xF1,0xFB), line=GRAY, line_w=1.0, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    fillh = Inches(3.6 * usedpct/100.0)
    box(s, x, Inches(2.3)+Inches(3.6)-fillh, Inches(1.4), fillh, fill=col)
    text(s, x, Inches(5.95), Inches(1.4), Inches(0.4), [(label, 11, DARK, True)], align=PP_ALIGN.CENTER)
    text(s, x, Inches(6.3), Inches(1.4), Inches(0.5), [(note, 10, GRAY, False)], align=PP_ALIGN.CENTER)
volbar(rx+Inches(0.4), "已用 85%", 85, GREEN, "未到阈值\n不变")
volbar(rx+Inches(2.0), "已用 90%", 90, ORANGE, "达 grow\n阈值→触发")
volbar(rx+Inches(3.6), "扩容后", 60, GREEN, "卷变大\n已用%降")
text(s, rx+Inches(0.45), Inches(2.0), Inches(5), Inches(0.3),
     [("虚线 = grow-threshold 90%", 10, ORANGE, True)])


# ============ Slide 3: Autosizing CLI ============
s = add_slide(); bg(s, WHITE)
title_bar(s, "Autosizing — CLI 配置", "volume autosize 命令（SSH 进 fsxadmin）")

text(s, Inches(0.7), Inches(1.4), Inches(12), Inches(0.4),
     [("① SSH 登录 ONTAP 管理端点", 16, NETAPP_BLUE, True)])
codebox(s, Inches(0.7), Inches(1.85), Inches(11.95), Inches(0.6),
        ["[~]$ ssh fsxadmin@<management_endpoint_ip>"])

text(s, Inches(0.7), Inches(2.65), Inches(12), Inches(0.4),
     [("② 启用 autosize（grow_shrink 模式：可增可减）", 16, NETAPP_BLUE, True)])
codebox(s, Inches(0.7), Inches(3.1), Inches(11.95), Inches(1.5),
        ["::> volume autosize -vserver <svm> -volume <vol> \\",
         "      -mode grow_shrink \\",
         "      -grow-threshold-percent 90 \\      # 已用90%触发扩容",
         "      -maximum-size 300TB \\             # 上限(最大300TB)",
         "      -shrink-threshold-percent 50 \\    # 已用降到50%触发缩容",
         "      -minimum-size 100GB                # 缩容下限"])

text(s, Inches(0.7), Inches(4.8), Inches(12), Inches(0.4),
     [("③ 查看当前 autosize 配置", 16, NETAPP_BLUE, True)])
codebox(s, Inches(0.7), Inches(5.25), Inches(11.95), Inches(0.6),
        ["::> volume autosize -vserver <svm> -volume <vol>"])

box(s, Inches(0.7), Inches(6.1), Inches(11.95), Inches(0.85), fill=LIGHT_BG, line=ORANGE, line_w=1.0, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
text(s, Inches(0.9), Inches(6.18), Inches(11.6), Inches(0.7),
     [("说明：mode 设 grow 则只增不减（无需 shrink 参数）；仅 FlexVol 支持；max-size 最大 300TB，默认=卷大小×120%。",
       12, ACCENT, True)], anchor=MSO_ANCHOR.MIDDLE)


# ============ Slide 4: Cloud Write Mode 原理 ============
s = add_slide(); bg(s, WHITE)
title_bar(s, "Cloud Write Mode 云直写 — 原理", "新写入直接落到容量池(S3)，跳过 SSD 主层")

lx = Inches(0.7)
box(s, lx, Inches(1.45), Inches(5.7), Inches(5.4), fill=LIGHT_BG, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
text(s, lx+Inches(0.3), Inches(1.6), Inches(5.1), Inches(0.5),
     [("是什么", 20, NETAPP_BLUE, True)])
text(s, lx+Inches(0.3), Inches(2.25), Inches(5.1), Inches(4.5),
     [("• 正常：新数据先写 SSD 主层，", 14, DARK, False),
      ("   冷却后再分层到容量池(S3)", 14, DARK, False, 2),
      ("• 开启云直写后：", 14, ORANGE, True, 10),
      ("   新写入数据直接落到容量池(S3)，", 14, DARK, False, 4),
      ("   绕过 SSD 主层", 14, DARK, False, 2),
      ("• 主要用途：大规模数据迁移", 14, GREEN, True, 10),
      ("   （如通过 NFS 灌入海量数据）", 12, GRAY, False, 2),
      ("• 好处：迁移数据不占满昂贵 SSD，", 14, DARK, False, 10),
      ("   直接进低成本弹性容量池", 14, DARK, False, 2),
      ("• 三个前提条件（缺一不可）：", 14, DARK, True, 10),
      ("   ① 必须是已存在的卷", 13, DARK, False, 4),
      ("   ② 必须是 RW（读写）卷", 13, DARK, False, 2),
      ("   ③ 分层策略必须为 All", 13, DARK, False, 2)])

# 右侧对比图
rx = Inches(6.7)
box(s, rx, Inches(1.45), Inches(5.95), Inches(5.4), fill=WHITE, line=NETAPP_BLUE, line_w=1.5, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
text(s, rx+Inches(0.25), Inches(1.55), Inches(5.4), Inches(0.4),
     [("写入路径对比", 16, NETAPP_BLUE, True)])
# 正常路径
text(s, rx+Inches(0.3), Inches(2.05), Inches(5.4), Inches(0.35),
     [("正常模式：", 13, DARK, True)])
box(s, rx+Inches(0.4), Inches(2.45), Inches(1.5), Inches(0.7), fill=GREEN, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
text(s, rx+Inches(0.4), Inches(2.45), Inches(1.5), Inches(0.7), [("写入", 12, WHITE, True)], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
text(s, rx+Inches(1.95), Inches(2.5), Inches(0.5), Inches(0.6), [("→", 18, GRAY, True)], anchor=MSO_ANCHOR.MIDDLE)
box(s, rx+Inches(2.45), Inches(2.45), Inches(1.4), Inches(0.7), fill=NETAPP_BLUE, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
text(s, rx+Inches(2.45), Inches(2.45), Inches(1.4), Inches(0.7), [("SSD 主层", 11, WHITE, True)], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
text(s, rx+Inches(3.9), Inches(2.5), Inches(0.5), Inches(0.6), [("→", 18, GRAY, True)], anchor=MSO_ANCHOR.MIDDLE)
box(s, rx+Inches(4.4), Inches(2.45), Inches(1.3), Inches(0.7), fill=CLOUD, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
text(s, rx+Inches(4.4), Inches(2.45), Inches(1.3), Inches(0.7), [("容量池S3", 11, WHITE, True)], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
text(s, rx+Inches(2.45), Inches(3.2), Inches(2.0), Inches(0.3), [("(冷却后分层)", 10, GRAY, False)], align=PP_ALIGN.CENTER)

# 云直写路径
text(s, rx+Inches(0.3), Inches(3.9), Inches(5.4), Inches(0.35),
     [("Cloud Write 模式：", 13, ORANGE, True)])
box(s, rx+Inches(0.4), Inches(4.3), Inches(1.5), Inches(0.7), fill=GREEN, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
text(s, rx+Inches(0.4), Inches(4.3), Inches(1.5), Inches(0.7), [("写入", 12, WHITE, True)], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
text(s, rx+Inches(1.95), Inches(4.35), Inches(2.2), Inches(0.6), [("———→ 直写", 14, ORANGE, True)], anchor=MSO_ANCHOR.MIDDLE)
box(s, rx+Inches(4.4), Inches(4.3), Inches(1.3), Inches(0.7), fill=CLOUD, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
text(s, rx+Inches(4.4), Inches(4.3), Inches(1.3), Inches(0.7), [("容量池S3", 11, WHITE, True)], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
# 灰掉的SSD
box(s, rx+Inches(2.45), Inches(4.3), Inches(1.4), Inches(0.7), fill=RGBColor(0xD9,0xDE,0xE4), shape=MSO_SHAPE.ROUNDED_RECTANGLE)
text(s, rx+Inches(2.45), Inches(4.3), Inches(1.4), Inches(0.7), [("跳过SSD", 11, GRAY, True)], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

box(s, rx+Inches(0.4), Inches(5.5), Inches(5.2), Inches(1.1), fill=LIGHT_BG, line=ORANGE, line_w=1.0, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
text(s, rx+Inches(0.55), Inches(5.6), Inches(4.9), Inches(0.95),
     [("迁移完成后建议关闭云直写，恢复正常写 SSD，", 11, ACCENT, True),
      ("以保证日常读写的低延迟性能。", 11, ACCENT, True, 2)])


# ============ Slide 5: Cloud Write Mode CLI ============
s = add_slide(); bg(s, WHITE)
title_bar(s, "Cloud Write Mode — CLI 配置", "需 advanced 模式 + volume modify")

text(s, Inches(0.7), Inches(1.4), Inches(12), Inches(0.4),
     [("① SSH 登录 + 进入 advanced 模式", 16, NETAPP_BLUE, True)])
codebox(s, Inches(0.7), Inches(1.85), Inches(11.95), Inches(1.0),
        ["[~]$ ssh fsxadmin@<management_endpoint_ip>",
         "",
         "FSx::> set -privilege advanced"])

text(s, Inches(0.7), Inches(3.05), Inches(12), Inches(0.4),
     [("② 开启 / 关闭云直写", 16, NETAPP_BLUE, True)])
codebox(s, Inches(0.7), Inches(3.5), Inches(11.95), Inches(1.3),
        ["# 开启 (true)",
         "FSx::> volume modify -vserver <svm> -volume <vol> \\",
         "         -is-cloud-write-enabled true",
         "",
         "# 关闭改为 false"])

box(s, Inches(0.7), Inches(5.0), Inches(11.95), Inches(1.85), fill=LIGHT_BG, line=ORANGE, line_w=1.0, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
text(s, Inches(0.95), Inches(5.15), Inches(11.5), Inches(1.6),
     [("⚠ 前提条件（官方，缺一不可）：", 14, ORANGE, True),
      ("① 只能对【已存在】的卷开启（不能建卷时设）", 13, DARK, False, 6),
      ("② 卷必须是 RW（读写）类型", 13, DARK, False, 4),
      ("③ 卷的分层策略必须为 All（tiering-policy = all）", 13, DARK, False, 4),
      ("用途：大规模数据迁移（NFS 灌入海量数据）；迁移完建议关闭。", 13, ACCENT, True, 6)])


# ============ Slide 6: Cloud Write 与存储效率（压缩/去重） ============
s = add_slide(); bg(s, WHITE)
title_bar(s, "Cloud Write 与存储效率", "云直写的数据会被压缩/去重吗？—— 会，但只走 inline")

lx = Inches(0.7)
box(s, lx, Inches(1.45), Inches(5.9), Inches(5.4), fill=LIGHT_BG, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
text(s, lx+Inches(0.3), Inches(1.6), Inches(5.3), Inches(0.5),
     [("结论", 20, NETAPP_BLUE, True)])
text(s, lx+Inches(0.3), Inches(2.25), Inches(5.3), Inches(4.5),
     [("• Cloud Write 本身不是压缩/去重功能，", 14, DARK, False),
      ("   它只决定数据「写到哪」", 13, GRAY, False, 2),
      ("• 直写容量池(S3)的数据【会】享受：", 14, ORANGE, True, 10),
      ("   ✓ inline 去重 (deduplication)", 13, GREEN, True, 4),
      ("   ✓ inline 压缩 (compression)", 13, GREEN, True, 2),
      ("   ✓ inline 压紧 (compaction)", 13, GREEN, True, 2),
      ("• 但【得不到】后台(background)去重/压缩：", 14, ORANGE, True, 10),
      ("   ✗ 后台效率任务只在 SSD 主层周期运行，", 13, RGBColor(0xC0,0x39,0x2B), True, 4),
      ("      云直写绕过了 SSD，所以拿不到", 12, GRAY, False, 2),
      ("• 存储效率作用在 4KiB 数据块层级，", 14, DARK, False, 10),
      ("   非文件层级", 12, GRAY, False, 2),
      ("• 分层到云时，已有的压缩/去重/压紧", 14, DARK, False, 8),
      ("   会被【保留】，省对象存储容量+传输费", 13, DARK, False, 2)])

# 右侧：inline vs background 示意
rx = Inches(6.8)
box(s, rx, Inches(1.45), Inches(5.85), Inches(5.4), fill=WHITE, line=NETAPP_BLUE, line_w=1.5, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
text(s, rx+Inches(0.25), Inches(1.55), Inches(5.3), Inches(0.4),
     [("两种存储效率，云直写只拿到一种", 15, NETAPP_BLUE, True)])

# inline 卡片（拿得到）
box(s, rx+Inches(0.35), Inches(2.15), Inches(5.15), Inches(1.75), fill=RGBColor(0xE7,0xF5,0xEC), line=GREEN, line_w=1.5, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
text(s, rx+Inches(0.55), Inches(2.28), Inches(4.8), Inches(1.55),
     [("✓ Inline（写盘前在内存中处理）", 14, GREEN, True),
      ("时机：数据写入路径上实时进行", 12, DARK, False, 4),
      ("云直写场景：直写 S3 的数据【会】经过", 12, DARK, False, 4),
      ("inline dedup + compression + compaction", 11, GREEN, True, 2)])

# background 卡片（拿不到）
box(s, rx+Inches(0.35), Inches(4.05), Inches(5.15), Inches(1.75), fill=RGBColor(0xFB,0xEC,0xEA), line=RGBColor(0xC0,0x39,0x2B), line_w=1.5, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
text(s, rx+Inches(0.55), Inches(4.18), Inches(4.8), Inches(1.55),
     [("✗ Background（写盘后周期任务）", 14, RGBColor(0xC0,0x39,0x2B), True),
      ("时机：数据落 SSD 后由 efficiency job 处理", 12, DARK, False, 4),
      ("云直写场景：绕过 SSD →【得不到】", 12, RGBColor(0xC0,0x39,0x2B), True, 4),
      ("后台 dedup / compression", 11, RGBColor(0xC0,0x39,0x2B), True, 2)])

text(s, rx+Inches(0.35), Inches(5.95), Inches(5.2), Inches(0.8),
     [("注：metadata 始终存 SSD 主层（即使数据直写云）。", 11, ACCENT, True),
      ("如需完整去重率，迁移后关闭云直写、让数据回 SSD 跑后台效率。", 11, ACCENT, True, 3)])


# ============ Slide 7: 总结对比 ============
s = add_slide(); bg(s, WHITE)
title_bar(s, "总结对比", "Autosizing vs Cloud Write Mode")

rows = [
    ("维度", "Autosizing 自动扩缩容", "Cloud Write Mode 云直写"),
    ("解决问题", "卷容量不够/浪费", "迁移占满 SSD 主层"),
    ("作用", "按已用% 自动 grow/shrink", "新写入直写容量池(S3)"),
    ("适用卷", "FlexVol", "已存在 + RW + All分层"),
    ("CLI 命令", "volume autosize", "volume modify\n-is-cloud-write-enabled"),
    ("是否需advanced", "否", "是 (set -privilege advanced)"),
    ("典型场景", "长期运行、容量波动", "一次性大规模迁移"),
]
tx = Inches(0.7); ty = Inches(1.45)
colw = [Inches(2.6), Inches(4.6), Inches(4.75)]
rh = Inches(0.7)
y = ty
for ri, row in enumerate(rows):
    x = tx
    header = ri == 0
    for ci, cell in enumerate(row):
        fill = NETAPP_BLUE if header else (LIGHT_BG if ri % 2 else WHITE)
        box(s, x, y, colw[ci], rh, fill=fill, line=RGBColor(0xD5,0xDE,0xEA), line_w=0.75)
        col = WHITE if header else (NETAPP_BLUE if ci==0 else (GREEN if ci==1 else CLOUD))
        sz = 14 if header else 12
        bold = header or ci==0
        text(s, x+Inches(0.12), y, colw[ci]-Inches(0.24), rh,
             [(cell, sz, col, bold)], anchor=MSO_ANCHOR.MIDDLE,
             align=PP_ALIGN.CENTER if (header or ci>=1) else PP_ALIGN.LEFT)
        x += colw[ci]
    y += rh

text(s, Inches(0.7), Inches(6.5), Inches(12), Inches(0.6),
     [("来源：AWS FSx for ONTAP 官方文档（autosizing / cloud write mode）+ NetApp FabricPool 文档 + AWS Storage Blog（Cloud Write inline 效率）。均通过 ONTAP CLI（SSH fsxadmin）配置。",
       12, ACCENT, True)])

prs.save("/home/ubuntu/.openclaw/workspace/FSx_ONTAP_Autosize_CloudWrite.pptx")
print("saved")
